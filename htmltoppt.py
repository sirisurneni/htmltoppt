from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import html
import sys
import os

def html_to_pptx(html_content, output_filename="presentation.pptx"):
    """
    Convert HTML slides to PowerPoint presentation using blank slides without placeholders
    
    Args:
        html_content (str): HTML content with slides
        output_filename (str): Output PowerPoint file name
    """
    # Create a new presentation
    prs = Presentation()
    
    # Parse HTML content
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Extract styles from the HTML
    css_rules = extract_css_rules(soup)
    
    # Find all slide divs
    slides = soup.find_all('div', class_='slide')
    
    for slide_index, slide in enumerate(slides):
        # Use a blank slide to avoid placeholders
        slide_layout = prs.slide_layouts[6]  # Blank slide
        current_slide = prs.slides.add_slide(slide_layout)
        
        # Add title manually instead of using placeholder
        title_element = slide.find('h1') or slide.find('h2')
        if title_element:
            title_shape = current_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(1)
            )
            title_frame = title_shape.text_frame
            p = title_frame.add_paragraph()
            p.text = title_element.text.strip()
            p.font.size = Pt(32)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        
        # Add a text box for the main content
        content_shape = current_slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
        content_frame = content_shape.text_frame
        
        # Process the slide content
        process_slide_content_without_placeholders(slide, content_frame, css_rules)
        
        # Clean up any lingering placeholders
        clean_slide_placeholders(current_slide)
    
    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")

def process_slide_content_without_placeholders(slide, text_frame, css_rules):
    """Process slide content without using PowerPoint placeholders"""
    # Find and process all row divs
    rows = slide.find_all('div', class_='row')
    
    # If no rows are found, process the slide content directly
    if not rows:
        process_content(slide, text_frame, css_rules)
    else:
        # Process each row
        for row in rows:
            # Create a paragraph separator
            p = text_frame.add_paragraph()
            
            # Apply any CSS styling from the row's class
            apply_css_to_paragraph(p, row, css_rules)
            
            # Process the content of the row
            process_content(row, text_frame, css_rules)
            
            # Add some spacing between rows
            if len(text_frame.paragraphs) > 0:
                last_p = text_frame.paragraphs[-1]
                if hasattr(last_p, 'space_after'):
                    last_p.space_after = Pt(12)

def extract_css_rules(soup):
    """Extract CSS rules from style tags in the HTML"""
    css_rules = {}
    
    # Find all style tags
    style_tags = soup.find_all('style')
    for style_tag in style_tags:
        style_content = style_tag.string
        if not style_content:
            continue
            
        # Extract class-based rules
        for rule in re.findall(r'\.([^\s{]+)\s*{([^}]+)}', style_content):
            class_name = rule[0]
            properties = {}
            
            # Extract properties
            for prop in re.findall(r'([^:;]+):\s*([^;]+);?', rule[1]):
                prop_name = prop[0].strip()
                prop_value = prop[1].strip()
                properties[prop_name] = prop_value
                
            css_rules[class_name] = properties
    
    return css_rules

def get_slide_placeholders(slide):
    """Get a mapping of placeholder names to placeholder objects"""
    placeholders = {}
    for placeholder in slide.placeholders:
        if hasattr(placeholder, 'name'):
            placeholders[placeholder.name.lower()] = placeholder
        elif hasattr(placeholder, 'placeholder_format') and hasattr(placeholder.placeholder_format, 'type'):
            # Use type as fallback
            placeholder_type = str(placeholder.placeholder_format.type)
            placeholders[placeholder_type] = placeholder
    
    return placeholders

def process_content(element, text_frame, css_rules):
    """Process HTML content and add it to a PowerPoint text frame"""
    
    # Handle different content types appropriately
    if element.find('table'):
        process_table(element.find('table'), text_frame, css_rules)
    elif element.find('ul') or element.find('ol'):
        process_list(element, text_frame, css_rules)
    elif element.find(['pre', 'code']) or element.find('div', class_='code-block'):
        process_code_block(element, text_frame, css_rules)
    elif element.find('img'):
        process_image(element, text_frame, css_rules)
    else:
        # Process text content
        process_text_content(element, text_frame, css_rules)

def process_text_content(element, text_frame, css_rules):
    """Process text content and add it to the text frame"""
    # Extract direct text content from the element (exclude nested elements)
    direct_text = ''
    for child in element.children:
        if isinstance(child, str):
            direct_text += child
            
    if direct_text.strip():
        p = text_frame.add_paragraph()
        p.text = direct_text.strip()
        apply_css_to_paragraph(p, element, css_rules)
    
    # Process paragraph elements
    paragraphs = element.find_all(['p', 'div', 'h3', 'h4'], recursive=False)
    for para in paragraphs:
        p = text_frame.add_paragraph()
        p.text = para.get_text().strip()
        apply_css_to_paragraph(p, para, css_rules)
        
        # Apply special formatting
        if para.name in ['h3', 'h4']:
            p.font.bold = True
            size_map = {'h3': 20, 'h4': 18}
            p.font.size = Pt(size_map.get(para.name, 16))
            
        # Handle text highlighting
        if para.find('span', class_='highlight'):
            # In a real implementation, you would extract the exact highlighted text
            # This is a simplification
            p.font.highlight_color = 3  # Yellow
            
        # Handle bold and italic
        if para.find(['b', 'strong']):
            p.font.bold = True
        if para.find(['i', 'em']):
            p.font.italic = True

def process_list(element, text_frame, css_rules):
    """Process HTML lists and add them to the text frame"""
    # First add any text before the list
    text_before = ''
    list_elem = element.find(['ul', 'ol'])
    
    for sibling in list_elem.previous_siblings:
        if isinstance(sibling, str) and sibling.strip():
            text_before += sibling.strip() + ' '
        elif hasattr(sibling, 'get_text'):
            text_before += sibling.get_text().strip() + ' '
            
    if text_before.strip():
        p = text_frame.add_paragraph()
        p.text = text_before.strip()
    
    # Process list items
    is_ordered = list_elem.name == 'ol'
    list_items = list_elem.find_all('li')
    
    for i, item in enumerate(list_items):
        p = text_frame.add_paragraph()
        prefix = f"{i+1}. " if is_ordered else "• "
        p.text = prefix + item.get_text().strip()
        p.level = 1  # Set indentation level
        
        apply_css_to_paragraph(p, item, css_rules)

def process_table(table, text_frame, css_rules):
    """Process HTML table and add it to the text frame as formatted text"""
    # Add table caption or heading
    p = text_frame.add_paragraph()
    p.text = "[Table]"
    p.font.bold = True
    
    # Process headers
    headers = [th.get_text().strip() for th in table.find_all('th')]
    if headers:
        p = text_frame.add_paragraph()
        p.text = " | ".join(headers)
        p.font.bold = True
        
        # Add separator line
        p = text_frame.add_paragraph()
        p.text = "-" * (sum(len(h) for h in headers) + 3 * (len(headers) - 1))
    
    # Process rows
    for row in table.find_all('tr'):
        cells = [td.get_text().strip() for td in row.find_all('td')]
        if cells:
            p = text_frame.add_paragraph()
            p.text = " | ".join(cells)

def process_code_block(element, text_frame, css_rules):
    """Process code blocks and add them to the text frame"""
    # Find the code block element
    code_elem = element.find(['pre', 'code']) or element.find('div', class_='code-block')
    
    if not code_elem:
        return
        
    # Add a label
    p = text_frame.add_paragraph()
    p.text = "[Code]"
    p.font.bold = True
    
    # Process code lines
    code_text = code_elem.get_text().strip()
    lines = code_text.split('\n')
    
    for line in lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(9)

def process_image(element, text_frame, css_rules):
    """Process images and add placeholders to the text frame"""
    img = element.find('img')
    if not img:
        return
        
    # Add an image placeholder
    p = text_frame.add_paragraph()
    img_alt = img.get('alt', 'Image')
    p.text = f"[Image: {img_alt}]"
    p.alignment = PP_ALIGN.CENTER
    
    # Add image caption if available
    caption = element.find('p', class_='caption')
    if caption:
        p = text_frame.add_paragraph()
        p.text = caption.get_text().strip()
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

def apply_css_to_paragraph(paragraph, element, css_rules):
    """Apply CSS styling to a PowerPoint paragraph based on element classes"""
    # Get classes from the element
    classes = element.get('class', [])
    if isinstance(classes, str):
        classes = classes.split()
        
    # Apply styling from each class
    for class_name in classes:
        if class_name in css_rules:
            props = css_rules[class_name]
            
            # Text alignment
            if 'text-align' in props:
                align_value = props['text-align'].lower()
                if align_value == 'center':
                    paragraph.alignment = PP_ALIGN.CENTER
                elif align_value == 'right':
                    paragraph.alignment = PP_ALIGN.RIGHT
                elif align_value == 'justify':
                    paragraph.alignment = PP_ALIGN.JUSTIFY
                    
            # Font size (approximate conversion from px/em to points)
            if 'font-size' in props:
                size_str = props['font-size']
                size_value = extract_numeric_value(size_str)
                
                if size_value:
                    # Convert common units to points (approximate)
                    if 'px' in size_str:
                        paragraph.font.size = Pt(size_value * 0.75)  # px to pt conversion
                    elif 'em' in size_str:
                        paragraph.font.size = Pt(size_value * 12)  # em to pt conversion
                    elif 'pt' in size_str:
                        paragraph.font.size = Pt(size_value)
                    else:
                        # Default unit or percentage
                        paragraph.font.size = Pt(size_value)
                        
            # Font weight
            if 'font-weight' in props:
                weight = props['font-weight'].lower()
                if weight in ['bold', 'bolder', '700', '800', '900']:
                    paragraph.font.bold = True
                    
            # Font style
            if 'font-style' in props:
                style = props['font-style'].lower()
                if style == 'italic':
                    paragraph.font.italic = True
                    
            # Text color (simplified conversion)
            if 'color' in props:
                color = props['color']
                rgb = extract_rgb_color(color)
                if rgb:
                    paragraph.font.color.rgb = RGBColor(*rgb)

def extract_numeric_value(value_str):
    """Extract numeric value from a CSS value string"""
    match = re.search(r'([0-9.]+)', value_str)
    if match:
        try:
            return float(match.group(1))
        except ValueError:
            pass
    return None

def extract_rgb_color(color_str):
    """Extract RGB values from a CSS color string"""
    # Handle hex colors
    hex_match = re.search(r'#([0-9a-fA-F]{6})', color_str)
    if hex_match:
        hex_value = hex_match.group(1)
        return (
            int(hex_value[0:2], 16),
            int(hex_value[2:4], 16),
            int(hex_value[4:6], 16)
        )
        
    # Handle rgb() format
    rgb_match = re.search(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
    if rgb_match:
        return (
            int(rgb_match.group(1)),
            int(rgb_match.group(2)),
            int(rgb_match.group(3))
        )
        
    return None

def clean_slide_placeholders(slide):
    """Remove or hide any empty placeholders on the slide"""
    for shape in slide.shapes:
        # Check if it's a placeholder
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            try:
                # Try setting it to empty to remove the "Click to add..." text
                if hasattr(shape, 'text'):
                    shape.text = ""
                # Or try to hide it
                if hasattr(shape, 'element') and hasattr(shape.element, 'getparent'):
                    parent = shape.element.getparent()
                    if parent is not None:
                        parent.remove(shape.element)
            except:
                # If we can't modify it, just continue
                pass

# Example usage
if __name__ == "__main__":
    # Default file names
    html_file = "sample1.html"          # Default input HTML file
    output_file = "presentation.ppt"   # Default output PowerPoint file
    
    # Check for command line arguments
    
    
    # Try to read the HTML file
    try:
        with open(html_file, 'r', encoding='utf-8') as f:
            html_content = f.read()
            
        # Convert HTML to PowerPoint
        html_to_pptx(html_content, output_file)
        print(f"Successfully converted {html_file} to {output_file}")
        
    except FileNotFoundError:
        print(f"File not found: {html_file}")
        print("Please ensure the HTML file exists or specify the correct path.")
        print("Usage: python html_to_pptx.py <html_file> [output_pptx]")
    except Exception as e:
        print(f"Error: {e}")
        print("Usage: python html_to_pptx.py <html_file> [output_pptx]")