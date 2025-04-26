from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
from io import BytesIO
from pptx.enum.text import MSO_AUTO_SIZE
from jinja2 import Template,Environment, FileSystemLoader
import json
import os
import sys
import requests
import re
import html
import sys
import os
import argparse
# Standard slide dimensions in inches
SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5
FOOTER_HEIGHT_INCHES = 0.5


def render_template_with_jinja(template_html, json_data):
    """
    Render the HTML template with JSON data using Jinja2's full capabilities
    
    Args:
        template_html (str): HTML template with Jinja2 placeholders
        json_data (dict): JSON data to fill the placeholders
        
    Returns:
        str: Rendered HTML with placeholders filled
    """
    # Create a Jinja2 Template object from the HTML string
    template = Template(template_html)
    
    # Render the template with the JSON data
    # This allows for full Jinja2 features like loops, conditionals, filters, etc.
    rendered_html = template.render(**json_data)
    
    return rendered_html
def add_banner_to_slide(slide, banner_url=None, title_height=Inches(1.4)):
    """
    Add a banner to the top of the slide - either from URL or default light blue
    
    Args:
        slide: The PowerPoint slide to add the banner to
        banner_url: URL of the banner image (optional)
        title_height: Height position where content starts (default 1.5 inches)
    
    Returns:
        None
    """
    # Banner dimensions
    banner_left = 0
    banner_top = 0
    banner_width = Inches(SLIDE_WIDTH_INCHES)  # Full slide width
    banner_height = title_height  # Height from top to where content starts
    
    if banner_url and banner_url.strip():
        # Try to download and use the banner image from URL
        try:
            # Download the image with timeout
            response = requests.get(banner_url, stream=True, timeout=15)
            
            if response.status_code == 200:
                # Create image from content
                img_bytes = BytesIO(response.content)
                
                try:
                    # Try to add the image as banner
                    banner = slide.shapes.add_picture(
                        img_bytes, 
                        banner_left, 
                        banner_top, 
                        width=banner_width,
                        height=banner_height
                    )
                    
                    print(f"Added banner from URL: {banner_url}")
                    
                    # Add Infosys text to top right corner
                    text_box = slide.shapes.add_textbox(
                        Inches(SLIDE_WIDTH_INCHES - 2), Inches(0.2),  # Position at top right
                        Inches(1.8), Inches(0.5)  # Size of text box
                    )
                    text_frame = text_box.text_frame
                    p = text_frame.add_paragraph()
                    p.text = "@INFOSYS"
                    p.font.bold = True
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue color for contrast
                    p.alignment = PP_ALIGN.RIGHT
                    
                    return  # Exit early as we successfully added the banner
                    
                except Exception as img_error:
                    print(f"Error adding banner image: {img_error}. Using default banner instead.")
                    # Continue to default banner creation below
                finally:
                    img_bytes.close()
            else:
                print(f"Failed to download banner image (status {response.status_code}). Using default banner.")
        except Exception as request_error:
            print(f"Error downloading banner image: {request_error}. Using default banner instead.")
    
    # If URL is not provided or any error occurred, create the default banner
    # Create banner shape - IMPORTANT: Add this FIRST before any other content
    # This ensures it's at the back without needing to modify slide structure later
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        banner_left, banner_top, banner_width, banner_height
    )
    
    # Set light blue color (RGB: 173, 216, 230) - standard light blue
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(173, 216, 230)
    
    # No border
    banner.line.width = 0
    
    # Add Infosys text to top right corner
    text_box = slide.shapes.add_textbox(
        Inches(SLIDE_WIDTH_INCHES - 2), Inches(0.2),  # Position at top right
        Inches(1.8), Inches(0.5)  # Size of text box
    )
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = "@INFOSYS"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue color for contrast
    p.alignment = PP_ALIGN.RIGHT

def render_template_file_with_jinja(template_path, json_data, template_dir=None):
    """
    Render a template file with JSON data using Jinja2's full capabilities
    including template inheritance, includes, etc.
    
    Args:
        template_path (str): Path to the template file (relative to template_dir)
        json_data (dict): JSON data to fill the placeholders
        template_dir (str): Directory containing template files
        
    Returns:
        str: Rendered HTML with placeholders filled
    """
    # If no template directory provided, use the current directory
    template_dir = template_dir or '.'
    
    # Create a Jinja2 Environment with the template directory
    env = Environment(loader=FileSystemLoader(template_dir))
    
    # Get the template from the environment
    template = env.get_template(template_path)
    
    # Render the template with the JSON data
    rendered_html = template.render(**json_data)
    
    return rendered_html
def generate_ppt_from_json_string_and_template_string(template_html, json_string, output_pptx="presentation.pptx"):
    """
    Generate a PowerPoint presentation from JSON string and HTML template string
    
    Args:
        template_html (str): HTML template string with Jinja2 placeholders
        json_string (str): JSON data as a string
        output_pptx (str): Path to save the PowerPoint file
        
    Returns:
        str: Path to the generated PowerPoint file
    """
    try:
        # Parse the JSON data
        json_data = json.loads(json_string)
        
        # Render the HTML template with the JSON data using full Jinja2 capabilities
        rendered_html = render_template_with_jinja(template_html, json_data)
        
        # Convert the rendered HTML to PowerPoint using your existing converter
        html_to_pptx(rendered_html, output_pptx)
        
        print(f"Generated PowerPoint presentation: {output_pptx}")
        return output_pptx
        
    except Exception as e:
        print(f"Error generating PowerPoint: {e}")
        raise
def generate_ppt_from_json_and_template(template_file, json_file, output_pptx="presentation.pptx", banner_url=None):
    """
    Generate a PowerPoint presentation from a JSON file and HTML template
    
    Args:
        template_file (str): Path to the HTML template file
        json_file (str): Path to the JSON data file
        output_pptx (str): Path to save the PowerPoint file
        banner_url (str): URL for the banner image (optional)
        
    Returns:
        str: Path to the generated PowerPoint file
    """
    try:
        # Get the template directory (the folder containing the template file)
        template_dir = os.path.dirname(template_file)
        template_name = os.path.basename(template_file)
        
        # Load the JSON data
        with open(json_file, 'r', encoding='utf-8') as f:
            json_data = json.load(f)
        
        # Use the file-based Jinja2 rendering
        rendered_html = render_template_file_with_jinja(template_name, json_data, template_dir)
        
        # Create a temporary HTML file with the rendered content
        temp_html_file = "temp_rendered.html"
        with open(temp_html_file, 'w', encoding='utf-8') as f:
            f.write(rendered_html)
        
        # Convert the rendered HTML to PowerPoint using your existing converter
        html_to_pptx(rendered_html, output_pptx, banner_url)
        
        # Optionally remove the temporary file
        # os.remove(temp_html_file)
        
        print(f"Generated PowerPoint presentation: {output_pptx}")
        return output_pptx
        
    except Exception as e:
        print(f"Error generating PowerPoint: {e}")
        raise
def html_to_pptx(html_content, output_filename="presentation.pptx", banner_url=None):
    """
    Convert HTML to PowerPoint presentation with support for mixed layouts
    
    Args:
        html_content (str): HTML content with slides
        output_filename (str): Output PowerPoint file name
        banner_url (str): URL for the banner image (optional)
    """
    # Create a new presentation
    prs = Presentation()
    
    # Parse HTML content
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Extract styles from the HTML
    
    
    # Find all slide divs
    slides = soup.find_all('div', class_='slide')
    
    # Process each slide based on its content
    for slide_index, slide_html in enumerate(slides):
        # Check if this slide has column layout
        left_column = slide_html.find('div', class_='left-column')
        right_column = slide_html.find('div', class_='right-column')
        use_columns_for_slide = bool(left_column or right_column)
        
        if use_columns_for_slide:
            # Process as column layout
            process_column_slide(slide_html, prs, slide_index, banner_url)
        else:
            # Process as standard layout
            process_standard_slide(slide_html, prs, slide_index, banner_url)
    
    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")

def process_standard_slide(slide, prs, slide_index, banner_url=None):
    """Process a slide with standard layout and apply background color if specified"""
    # Use a blank slide to avoid placeholders
    slide_layout = prs.slide_layouts[6]  # Blank slide
    current_slide = prs.slides.add_slide(slide_layout)
    
    # First add the banner - MUST be first to ensure it's at the back
    add_banner_to_slide(current_slide, banner_url, Inches(1.4))
    
    # Apply background color if the slide has a color class
    apply_slide_background_color(slide, current_slide)
    
    # Add title manually instead of using placeholder
    title_element = slide.find('h1') or slide.find('h2')
    
    if title_element:
        title_shape = current_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_frame = title_shape.text_frame
        p = title_frame.add_paragraph()
        p.text = title_element.text.strip()
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    # Process the slide content - now passing prs and slide_index
    process_standard_slide_content(slide, current_slide, prs, slide_index)
    add_footer(current_slide)
    # Clean up any lingering placeholders
    clean_slide_placeholders(current_slide)




def get_color_from_class(element, default_color=RGBColor(255, 255, 255)):
    """Extract background color based on color classes (red, blue, green, etc.)"""
    # Standard color mapping
    color_map = {
        'red': RGBColor(255, 200, 200),     # Light red
        'blue': RGBColor(200, 200, 255),    # Light blue
        'green': RGBColor(200, 255, 200),   # Light green
        'yellow': RGBColor(255, 255, 200),  # Light yellow
        'orange': RGBColor(255, 225, 180),  # Light orange
        'purple': RGBColor(230, 200, 255),  # Light purple
        'grey': RGBColor(220, 220, 220),    # Light grey
        'gray': RGBColor(220, 220, 220),    # Light gray
        'pink': RGBColor(255, 200, 230),    # Light pink
        'teal': RGBColor(180, 240, 240),    # Light teal
    }
    
    # Check if element has any of the color classes
    classes = element.get('class', [])
    if isinstance(classes, str):
        classes = classes.split()
        
    for cls in classes:
        if cls.lower() in color_map:
            return color_map[cls.lower()]
            
    # Return default if no color class found
    return default_color


# Also update the handle_text_overflow function to manage text better


# Modified text handling functions to properly wrap text and prevent slide overflow

def process_text_content(element, text_frame, css_rules, slide=None, prs=None, slide_index=0):
    """Process text content and add it to the text frame with improved text wrapping"""
    # Enable word wrap for the text frame
    text_frame.word_wrap = True
    
    # Set appropriate text frame margins
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    
    # Extract all text with a simpler approach
    all_text = element.get_text().strip()
    
    # If no text, return early
    if not all_text:
        return
    
    # Calculate how much text might fit based on the text frame dimensions
    # This is a simplified estimate - in practice, PowerPoint handles wrapping
    if slide and prs and len(all_text) > 800:  # Reduced from 1000 for better fit
        # Use the text overflow handler for long text
        handle_text_overflow(all_text, text_frame, slide, slide_index, prs)
    else:
        # Use smart paragraph splitting for better text flow
        paragraphs = all_text.split('\n')
        for para_text in paragraphs:
            if not para_text.strip():
                continue
                
            p = text_frame.add_paragraph()
            p.text = para_text.strip()
            
            # Apply basic formatting if needed
            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                p.font.bold = True
                size_map = {'h1': 28, 'h2': 24, 'h3': 24, 'h4': 18, 'h5': 16, 'h6': 14}
                p.font.size = Pt(size_map.get(element.name, 14))
            elif element.name in ['strong', 'b']:
                p.font.bold = True
            elif element.name in ['em', 'i']:
                p.font.italic = True


def handle_text_overflow(text, text_frame, slide, current_slide_index, prs):
    """Break long text content across multiple slides with improved text wrapping"""
    # Use a more conservative character count to ensure text fits
    chars_per_slide = 600  # Even more conservative than before
    
    # Split text into paragraphs first for better formatting
    paragraphs = text.split('\n')
    
    current_chars = 0
    current_para_index = 0
    
    # Add paragraphs until we hit the character limit
    while current_para_index < len(paragraphs):
        para_text = paragraphs[current_para_index].strip()
        
        # Skip empty paragraphs
        if not para_text:
            current_para_index += 1
            continue
        
        # If adding this paragraph would exceed our limit, create a continuation slide
        if current_chars + len(para_text) > chars_per_slide and current_chars > 0:
            # We need to continue on a new slide
            next_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            
            # First add the banner - MUST be first to ensure proper layering
            add_banner_to_slide(slide, banner_url, Inches(1.5))
            
            # Add a title indicating continuation
            title_shape = next_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
            )
            title_frame = title_shape.text_frame
            p = title_frame.add_paragraph()
            p.text = f"Continued from Slide {current_slide_index+1}"
            p.font.italic = True
            p.font.bold = True
            p.font.size = Pt(18)
            
            # Add the content with better positioning
            next_text_shape = next_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
            )
            next_text_frame = next_text_shape.text_frame
            next_text_frame.word_wrap = True
            next_text_frame.margin_left = 0
            next_text_frame.margin_right = 0
            next_text_frame.margin_top = 0
            next_text_frame.margin_bottom = 0
            
            # Recursively handle remaining paragraphs
            remaining_paras = paragraphs[current_para_index:]
            remaining_text = '\n'.join(remaining_paras)
            handle_text_overflow(remaining_text, next_text_frame, next_slide, 
                               current_slide_index+1, prs)
            return True
        
        # If we get here, we can add this paragraph to the current slide
        p = text_frame.add_paragraph()
        p.text = para_text
        
        current_chars += len(para_text)
        current_para_index += 1
    
    return False

def process_standard_slide_content(slide_html, current_slide, prs=None, slide_index=0):
    """Process content for a standard slide layout with better content fitting"""
    # Track vertical position for adding content
    current_y = Inches(1.5)  # Start after title
    
    # Calculate maximum content height
    max_y = Inches(SLIDE_HEIGHT_INCHES - 0.7 - FOOTER_HEIGHT_INCHES)

    
    # Get overall text length to determine if we need overflow handling
    full_text = slide_html.get_text().strip()
    
    # If the entire content is very long, handle it specially
    if len(full_text) > 1000 and prs:  # Lower threshold for better content fit
        content_shape = current_slide.shapes.add_textbox(
            Inches(0.5), current_y, Inches(9), Inches(5)
        )
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Handle as overflow text
        handle_text_overflow(full_text, content_frame, current_slide, slide_index, prs)
        return
    
    # Find and process all row divs
    rows = slide_html.find_all('div', class_='row')
    
    # If no rows are found, process the slide content directly
    if not rows:
        content_shape = current_slide.shapes.add_textbox(
            Inches(0.5), current_y, Inches(9), Inches(5)
        )
        content_frame = content_shape.text_frame
        process_content(slide_html, content_frame, current_slide, current_y, prs, slide_index)
    else:
        # Process each row with better spacing management
        for i, row in enumerate(rows):
            # Check remaining space
            remaining_height = max_y - current_y
            if remaining_height < Inches(1.0) and i < len(rows) - 1:
                # Not enough space for meaningful content
                # Create a new slide for remaining content
                if prs:
                    next_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # First add the banner - MUST be first to ensure proper layering
                    add_banner_to_slide(slide_html, banner_url, Inches(1.5))
                    # Add a title indicating continuation
                    title_element = slide_html.find('h1') or slide_html.find('h2')
                    title_text = title_element.get_text().strip() if title_element else f"Slide {slide_index+1}"
                    
                    title_shape = next_slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
                    )
                    title_frame = title_shape.text_frame
                    p = title_frame.add_paragraph()
                    p.text = f"{title_text} (Continued)"
                    p.font.italic = True
                    p.font.bold = True
                    p.font.size = Pt(18)
                    
                    # Process remaining rows on new slide
                    next_y = Inches(1.5)
                    for next_row in rows[i:]:
                        # Calculate content height
                        row_height = estimate_row_height(next_row)
                        
                        # Check if it fits on the continuation slide
                        if next_y + row_height > Inches(SLIDE_HEIGHT_INCHES - 0.7):
                            # Still too much content, need another slide
                            continue_index = rows.index(next_row)
                            if continue_index < len(rows) - 1:
                                # Recursively handle remaining content
                                remaining_rows_html = BeautifulSoup('<div></div>', 'html.parser').div
                                for r in rows[continue_index:]:
                                    remaining_rows_html.append(r.copy())
                                
                                process_standard_slide_content(
                                    remaining_rows_html, next_slide, prs, slide_index+1
                                )
                                break
                        
                        # Create a text frame for this row
                        text_shape = next_slide.shapes.add_textbox(
                            Inches(0.5), next_y, Inches(9), row_height
                        )
                        text_frame = text_shape.text_frame
                        text_frame.word_wrap = True
                        
                        # Process the content of the row
                        new_y = process_content(next_row, text_frame, next_slide, 
                                             next_y, prs, slide_index+1)
                        
                        # Update position for next row
                        next_y = max(next_y + row_height, new_y) + Inches(0.3) if new_y else next_y + row_height + Inches(0.3)
                    
                    # No need to process more rows on the original slide
                    break
            
            # Estimate row height with more conservative calculation
            row_height = estimate_row_height(row)
            
            # Adjust height if remaining space is limited
            if current_y + row_height > max_y:
                row_height = max_y - current_y - Inches(0.1)
            
            # Create a text frame for this row
            text_shape = current_slide.shapes.add_textbox(
                Inches(0.5), current_y, Inches(9), row_height
            )
            text_frame = text_shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            # Process the content of the row
            new_y = process_content(row, text_frame, current_slide, current_y, prs, slide_index)
            
            # Update the vertical position for the next row
            current_y = max(current_y + row_height, new_y) + Inches(0.2) if new_y else current_y + row_height + Inches(0.2)
            
            # Check if we're running out of space
            if current_y >= max_y and i < len(rows) - 1:
                # Create a new slide for remaining content
                if prs:
                    next_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # First add the banner - MUST be first to ensure proper layering
                    add_banner_to_slide(slide_html, banner_url, Inches(1.5))
                    # Add a title indicating continuation
                    title_element = slide_html.find('h1') or slide_html.find('h2')
                    title_text = title_element.get_text().strip() if title_element else f"Slide {slide_index+1}"
                    
                    title_shape = next_slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
                    )
                    title_frame = title_shape.text_frame
                    p = title_frame.add_paragraph()
                    p.text = f"{title_text} (Continued)"
                    p.font.italic = True
                    p.font.bold = True
                    p.font.size = Pt(18)
                    
                    # Recursively process remaining rows on new slide
                    remaining_rows_html = BeautifulSoup('<div></div>', 'html.parser').div
                    for r in rows[i+1:]:
                        remaining_rows_html.append(r.copy())
                    
                    process_standard_slide_content(
                        remaining_rows_html, next_slide, prs, slide_index+1
                    )
                break


# Two specific fixes for the HTML to PowerPoint converter:
# 1. Better processing of colors in div class tags for headings and paragraphs
# 2. Fix for the image in slide 2's right column to keep it inside the row box

# FIX 1: Improved color handling from div tags
def process_headers_with_color(element, text_frame):
    """Process headers with improved color styling"""
    for header in element.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
        p = text_frame.add_paragraph()
        p.text = header.get_text().strip()
        p.font.bold = True
        size_map = {'h1': 24, 'h2': 20, 'h3': 20, 'h4': 16, 'h5': 14, 'h6': 12}
        p.font.size = Pt(size_map.get(header.name, 14))
        
        # Check color in this priority: header tag itself, parent div, grandparent div
        header_color = get_color_from_class(header)
        
        # If no color found in the header itself, look at parent div
        if header_color == RGBColor(255, 255, 255):
            parent_div = header.find_parent('div')
            if parent_div:
                header_color = get_color_from_class(parent_div)
                
                # If still no color, try grandparent
                if header_color == RGBColor(255, 255, 255):
                    grandparent_div = parent_div.find_parent('div')
                    if grandparent_div:
                        header_color = get_color_from_class(grandparent_div)
        
        # Apply the color if one was found
        if header_color != RGBColor(255, 255, 255):
            p.font.color.rgb = header_color


def process_paragraphs_with_color(element, text_frame):
    """Process paragraphs with improved color styling"""
    for para in element.find_all('p'):
        p = text_frame.add_paragraph()
        
        # Get the text and highlight numbers with regex
        text = para.get_text().strip()
        
        # Find all numbers in the text
        num_positions = [(m.start(), m.end()) for m in re.finditer(r'\b\d+(\.\d+)?\b', text)]
        
        if num_positions:
            # There are numbers in this text
            last_pos = 0
            for start, end in num_positions:
                # Add text before the number
                if start > last_pos:
                    run = p.add_run()
                    run.text = text[last_pos:start]
                    run.font.size = Pt(12)
                
                # Add the number with bold formatting
                run = p.add_run()
                run.text = text[start:end]
                run.font.bold = True
                run.font.size = Pt(14)  # Slightly larger
                
                last_pos = end
            
            # Add any remaining text after the last number
            if last_pos < len(text):
                run = p.add_run()
                run.text = text[last_pos:]
                run.font.size = Pt(12)
        else:
            # No numbers, just add the text
            p.text = text
            p.font.size = Pt(12)
        
        # Apply color as before
        para_color = get_color_from_class(para)
        if para_color == RGBColor(255, 255, 255):
            parent_div = para.find_parent('div')
            if parent_div:
                para_color = get_color_from_class(parent_div)
                if para_color == RGBColor(255, 255, 255):
                    grandparent_div = parent_div.find_parent('div')
                    if grandparent_div:
                        para_color = get_color_from_class(grandparent_div)
        
        if para_color != RGBColor(255, 255, 255):
            p.font.color.rgb = para_color

# FIX 2: Keep images inside row boxes in column layouts
# Targeted fix for image overlap in column content while keeping everything in the same box


def process_column_slide(slide_html, prs, slide_idx,banner_url=None):
    """Process a slide with column layout and apply background color if specified"""
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # First add the banner - MUST be first to ensure it's at the back
    add_banner_to_slide(slide, banner_url, Inches(1.4))
    
    # Apply background color if the slide has a color class
    apply_slide_background_color(slide_html, slide)

    # Title
    title_element = slide_html.find('h1') or slide_html.find('h2')
    title_text = title_element.get_text().strip() if title_element else f"Slide {slide_idx + 1}"

    # Use standard slide dimensions
    slide_width_inches = SLIDE_WIDTH_INCHES
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(slide_width_inches - 1), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    # Left and Right columns
    left_column = slide_html.find('div', class_='left-column')
    right_column = slide_html.find('div', class_='right-column')

    # Column layout setup
    # Calculate dynamic column widths
    margin = Inches(0.5)
    col_spacing = Inches(0.5)
    
    usable_width = Inches(slide_width_inches - 1 - 0.5)  # Total width minus margins
    col_width = (usable_width - col_spacing) / 2  # Equal width for both columns
    
    left_x = margin
    right_x = margin + col_width + col_spacing
    start_y = Inches(1.5)  # Start below title

    # Store the final Y position after processing columns
    final_y_positions = []

    # Process left column if it exists
    y_left = start_y
    if left_column:
        print(f"Processing left column with {len(left_column.find_all('div', class_='row'))} rows")
        y_left = process_column_content(left_column, slide, left_x, y_left, col_width, slide_idx, prs)
        final_y_positions.append(y_left)

    # Process right column if it exists
    y_right = start_y
    if right_column:
        print(f"Processing right column with {len(right_column.find_all('div', class_='row'))} rows")
        y_right = process_column_content(right_column, slide, right_x, y_right, col_width, slide_idx, prs)
        final_y_positions.append(y_right)

    # Determine the highest Y position after processing both columns
    if final_y_positions:
        highest_y = max(final_y_positions) + Inches(0.3)  # Add margin after columns
    else:
        highest_y = start_y + Inches(0.5)  # Default if no columns were processed

    # IMPROVED: Find only rows that are direct children of the slide div or have the standalone class
    standalone_rows = []
    
    # First look specifically for standalone class
    for row in slide_html.find_all('div', class_='standalone'):
        standalone_rows.append(row)
        
    # Check for rows that are direct children of the slide div - not in any column
    for row in slide_html.find_all('div', class_='row', recursive=False):
        if row not in standalone_rows:
            # Double-check this row is not already in a column
            if not (row.parent.get('class') and ('left-column' in row.parent.get('class') or 'right-column' in row.parent.get('class'))):
                standalone_rows.append(row)
    
    print(f"Found {len(standalone_rows)} standalone rows")
    
    # IMPROVED: Calculate available space on slide and manage continuation
    max_slide_y = Inches(SLIDE_HEIGHT_INCHES - 0.7 - FOOTER_HEIGHT_INCHES)

    current_slide = slide
    
    # Process each standalone row with better space management
    for row_index, row in enumerate(standalone_rows):
        # Check for space and create continuation slide if needed
        if highest_y > max_slide_y and row_index > 0:
            # Need to create a continuation slide
            continuation_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # First add the banner - MUST be first for proper layering
            add_banner_to_slide(continuation_slide, banner_url,Inches(1.4))
            
            # Add continuation title
            cont_title_box = continuation_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(slide_width_inches - 1), Inches(1)
            )
            cont_title_frame = cont_title_box.text_frame
            cont_title_text = f"{title_text} (Continued)"
            cont_title_frame.text = cont_title_text
            cont_title_frame.paragraphs[0].font.size = Pt(28)
            cont_title_frame.paragraphs[0].font.bold = True
            
            # Reset positioning for new slide
            current_slide = continuation_slide
            highest_y = Inches(1.5)  # Start below title
            
            # Apply background color to continuation slide if needed
            apply_slide_background_color(slide_html, current_slide)
            
            # Clean up placeholders on the new slide
            clean_slide_placeholders(current_slide)
        
        # Full width for standalone rows
        full_width = Inches(slide_width_inches - 1)
        
        # Process the row on the current slide
        row_height = process_standalone_row(row, current_slide, margin, highest_y, full_width, slide_idx, prs)
        
        # Update the highest Y position for next row
        highest_y = row_height + Inches(0.2)  # Add spacing between rows
    add_footer(current_slide)
    # Clean up any lingering placeholders on the original slide
    clean_slide_placeholders(slide)
def process_standalone_row(row, slide, left_x, y_pos, width, slide_index, prs):
    """Process rows that appear below columns, spanning the full width"""
    try:
        print(f"Processing standalone row with content: {row.get_text().strip()[:50]}...")
        
        # Extract content from the row
        img_tags = row.find_all('img')
        has_images = len(img_tags) > 0
        
        # Extract text content
        header_text = ""
        for header in row.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            header_text += header.get_text().strip() + " "
        
        paragraph_text = ""
        for para in row.find_all('p'):
            paragraph_text += para.get_text().strip() + " "
        
        # Extract any other text
        other_text = ""
        row_copy = BeautifulSoup(str(row), 'html.parser')
        for tag in row_copy.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'img']):
            tag.decompose()
        
        for element in row_copy.descendants:
            if isinstance(element, str) and element.strip():
                other_text += element.strip() + " "
        
        combined_text = (header_text + " " + paragraph_text + " " + other_text).strip()
        has_text = bool(combined_text)
        
        # Get row background color
        row_color = get_color_from_class(row)
        
        # Calculate box height based on content
        text_length = len(combined_text)
        
        # Estimate text height based on length
        if text_length < 100:
            text_height = Inches(0.6)
        elif text_length < 250:
            text_height = Inches(1.0)
        elif text_length < 500:
            text_height = Inches(1.5)
        else:
            text_height = Inches(2.0)
        
        # Calculate image height if present
        image_height = Inches(0)
        if has_images:
            img = img_tags[0]
            if img.get('height'):
                try:
                    img_height = int(img.get('height'))
                    image_height = Inches(img_height / 96 + 0.4)
                except (ValueError, TypeError):
                    image_height = Inches(1.5)  # Default if parsing fails
            else:
                image_height = Inches(1.5)  # Default image height
        
        # Calculate total box height with text and image
        # Calculate total box height using the new function
        box_height = calculate_dynamic_box_height(header_text, paragraph_text, other_text, has_images, image_height)
        
        # Create background box
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            left_x, y_pos, 
            width, box_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = row_color
        bg_shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Create text box
        text_box = slide.shapes.add_textbox(
            left_x + Inches(0.1),
            y_pos + Inches(0.1),
            width - Inches(0.2),
            text_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        
        # Add header text
        if header_text.strip():
            p = text_frame.add_paragraph()
            p.text = header_text.strip()
            p.font.bold = True
            p.font.size = Pt(14)
            p.space_before = 0
            p.space_after = Pt(2)
        
        # Add paragraph text with number highlighting
        if paragraph_text.strip():
            p = text_frame.add_paragraph()
            
            # Find all numbers in the text
            num_positions = [(m.start(), m.end()) for m in 
                            re.finditer(r'\b\d+(\.\d+)?\b', paragraph_text.strip())]
            
            if num_positions:
                # There are numbers in this text - use runs to format differently
                last_pos = 0
                for start, end in num_positions:
                    # Add text before the number
                    if start > last_pos:
                        before_text = paragraph_text.strip()[last_pos:start]
                        before_run = p.add_run()
                        before_run.text = before_text
                        before_run.font.size = Pt(12)
                        before_run.font.bold = False
                    
                    # Add the number with bold formatting and larger font
                    num_text = paragraph_text.strip()[start:end]
                    num_run = p.add_run()
                    num_run.text = num_text
                    num_run.font.bold = True
                    num_run.font.size = Pt(14)  # Slightly larger for numbers
                    
                    last_pos = end
                
                # Add any remaining text after the last number
                if last_pos < len(paragraph_text.strip()):
                    after_text = paragraph_text.strip()[last_pos:]
                    after_run = p.add_run()
                    after_run.text = after_text
                    after_run.font.size = Pt(12)
                    after_run.font.bold = False
            else:
                # No numbers, just add the text normally
                p.text = paragraph_text.strip()
                p.font.bold = False
                p.font.size = Pt(12)
                
            p.space_before = 0
            p.space_after = 0
        
        # Add other text if present
        if other_text.strip():
            p = text_frame.add_paragraph()
            p.text = other_text.strip()
            p.font.bold = False
            p.font.size = Pt(12)
        
        # Process images if present
        if has_images:
            # Calculate image position - below text
            img_y = y_pos + text_height + Inches(0.3)
            
            for img in img_tags:
                try:
                    img_url = img.get('src', '')
                    
                    if img_url:
                        response = requests.get(img_url, stream=True, timeout=15)  # Increased timeout
                        if response.status_code == 200:
                            img_bytes = BytesIO(response.content)
                            
                            try:
                                # Get dimensions from image
                                with PILImage.open(img_bytes) as pil_img:
                                    aspect_ratio = pil_img.width / pil_img.height
                                    
                                    img_bytes.seek(0)  # Reset file pointer
                                    
                                    # Calculate image size
                                    if img.get('width') and img.get('height'):
                                        try:
                                            width_px = int(img.get('width'))
                                            height_px = int(img.get('height'))
                                            img_width = Inches(width_px / 96)
                                            img_height = Inches(height_px / 96)
                                        except (ValueError, TypeError):
                                            img_width = min(Inches(3.0), width - Inches(0.4))
                                            img_height = img_width / aspect_ratio
                                    else:
                                        img_width = min(Inches(3.0), width - Inches(0.4))
                                        img_height = img_width / aspect_ratio
                                    
                                    # Ensure image fits within width
                                    if img_width > width - Inches(0.4):
                                        img_width = width - Inches(0.4)
                                        img_height = img_width / aspect_ratio
                                    
                                    # Center the image
                                    img_x = left_x + (width - img_width) / 2
                                    
                                    # Make sure image doesn't exceed box height
                                    if (img_y + img_height) > (y_pos + box_height - Inches(0.1)):
                                        img_height = y_pos + box_height - img_y - Inches(0.1)
                                        img_width = img_height * aspect_ratio
                                    
                                    # Create a fresh copy of the image data
                                    img_data = BytesIO(img_bytes.getvalue())
                                    
                                    if img_height > Inches(0.2):  # Only add if reasonable size
                                        picture = slide.shapes.add_picture(
                                            img_data, 
                                            img_x, 
                                            img_y, 
                                            width=img_width, 
                                            height=img_height
                                        )
                                        print(f"Added image from {img_url}")
                                    
                                    # Close the copy
                                    img_data.close()
                            except Exception as img_error:
                                print(f"Error processing image: {img_error}")
                            
                            img_bytes.close()
                except Exception as img_error:
                    print(f"Error with image: {img_error}")
        
        # Return the Y position after this row
        return y_pos + box_height
        
    except Exception as row_error:
        print(f"Error processing standalone row: {row_error}")
        return y_pos + Inches(0.5)  # Default return if error occurs




from pptx.enum.text import MSO_AUTO_SIZE

def process_column_content(column, slide, x_pos, y_pos, width, slide_index=0, prs=None):
    """Process content of a column with proper handling of rows and images"""
    current_y = y_pos
    
    try:
        # Find all rows in this column
        rows = column.find_all('div', class_='row')
        print(f"Processing column content with {len(rows)} rows")
        
        # Process each row in the column
        for row in rows:
            try:
                # Check remaining space on slide
                #remaining_height = Inches(SLIDE_HEIGHT_INCHES - 1.0 - FOOTER_HEIGHT_INCHES) - current_y
                remaining_height = Inches(SLIDE_HEIGHT_INCHES - 1.0 - FOOTER_HEIGHT_INCHES) - current_y

                
                # If not enough space, create continuation slide
                if remaining_height < Inches(0.5) and prs:
                    # Create continuation slide logic
                    slide_layout = prs.slide_layouts[6]
                    next_slide = prs.slides.add_slide(slide_layout)
                    
                    # First add the banner - MUST be first to ensure proper layering
                    add_banner_to_slide(next_slide, banner_url, Inches(1.4))
                    
                    # Add continuation title
                    title_element = column.parent.find('h1') or column.parent.find('h2')
                    title_text = title_element.get_text().strip() if title_element else f"Slide {slide_index + 1}"
                    
                    title_box = next_slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(SLIDE_WIDTH_INCHES - 1), Inches(0.8)
                    )
                    title_frame = title_box.text_frame
                    p = title_frame.add_paragraph()
                    p.text = f"{title_text} (Continued)"
                    p.font.size = Pt(24)
                    p.font.bold = True
                    
                    # Process remaining rows on new slide
                    remaining_rows_index = rows.index(row)
                    if remaining_rows_index >= 0:
                        new_column = BeautifulSoup("<div></div>", "html.parser").div
                        for r in rows[remaining_rows_index:]:
                            new_column.append(r.copy())
                        
                        process_column_content(new_column, next_slide, x_pos, Inches(1.5), width, slide_index + 1, prs)
                    break
                
                # Extract content from this row
                img_tags = row.find_all('img')
                has_images = len(img_tags) > 0
                
                # Get text content
                header_text = ""
                for header in row.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                    header_text += header.get_text().strip() + " "
                
                paragraph_text = ""
                for para in row.find_all('p'):
                    paragraph_text += para.get_text().strip() + " "
                
                # Get any other text
                other_text = ""
                row_copy = BeautifulSoup(str(row), 'html.parser')
                for tag in row_copy.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'img']):
                    tag.decompose()
                
                for element in row_copy.descendants:
                    if isinstance(element, str) and element.strip():
                        other_text += element.strip() + " "
                
                # Combine all text
                combined_text = (header_text + " " + paragraph_text + " " + other_text).strip()
                has_text = bool(combined_text)
                
                # Get background color
                row_color = get_color_from_class(row)
                
                # Calculate space needed based on content
                text_length = len(combined_text)
                
                # Estimate text height
                if text_length < 100:
                    text_height = Inches(0.6)
                elif text_length < 250:
                    text_height = Inches(1.0)
                elif text_length < 500:
                    text_height = Inches(1.5)
                else:
                    text_height = Inches(2.0)
                
                # Add height for images if present
                image_height = Inches(0)
                if has_images:
                    # IMPROVEMENT: Add more space for images
                    img = img_tags[0]
                    if img.get('height'):
                        try:
                            img_height = int(img.get('height'))
                            # Increase the multiplier to allow more space
                            image_height = Inches((img_height / 96) * 1.5)
                        except (ValueError, TypeError):
                            image_height = Inches(1.5)  # Increased default
                    else:
                        image_height = Inches(1.5)  # Increased default
                
                # Calculate buffer space for images - INCREASED
                buffer_space = Inches(0.5) if has_images else Inches(0)
                
                # Calculate total box height with better sizing
                if has_images:
                    # IMPROVEMENT: Add more padding for images
                    box_height = text_height + buffer_space + image_height + Inches(0.6)
                else:
                    box_height = text_height + Inches(0.2)
                
                # Use dynamic height calculation if available
                try:
                    dynamic_height = calculate_dynamic_box_height(header_text, paragraph_text, other_text, has_images, image_height)
                    box_height = max(box_height, dynamic_height)  # Use whichever is larger for safety
                except:
                    # If calculate_dynamic_box_height is not defined or fails, use the regular calculation
                    pass
                
                # IMPROVEMENT: Add minimum height guarantee for image-containing rows
                if has_images:
                    min_height_with_image = text_height + buffer_space + image_height + Inches(0.8)
                    box_height = max(box_height, min_height_with_image)
                
                # Create background shape with enough height for all content
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    x_pos, current_y, 
                    width, box_height
                )
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = row_color
                bg_shape.line.color.rgb = RGBColor(200, 200, 200)
                
                # IMPROVEMENT: Separate text box and image placement
                # Create text box for content with reduced height
                text_height_actual = min(text_height, box_height - Inches(0.4))
                text_box = slide.shapes.add_textbox(
                    x_pos + Inches(0.1),
                    current_y + Inches(0.1),
                    width - Inches(0.2),
                    text_height_actual
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_top = 0
                text_frame.margin_bottom = 0
                text_frame.margin_left = 0
                text_frame.margin_right = 0
                
                # Add header text
                if header_text.strip():
                    p = text_frame.add_paragraph()
                    p.text = header_text.strip()
                    p.font.bold = True
                    p.font.size = Pt(14)
                    p.space_before = 0
                    p.space_after = Pt(2)
                
                # Add paragraph text with number highlighting
                if paragraph_text.strip():
                    p = text_frame.add_paragraph()
                    
                    # Find all numbers in the text
                    num_positions = [(m.start(), m.end()) for m in 
                                    re.finditer(r'\b\d+(\.\d+)?\b', paragraph_text.strip())]
                    
                    if num_positions:
                        # There are numbers in this text - use runs to format differently
                        last_pos = 0
                        for start, end in num_positions:
                            # Add text before the number
                            if start > last_pos:
                                before_text = paragraph_text.strip()[last_pos:start]
                                before_run = p.add_run()
                                before_run.text = before_text
                                before_run.font.size = Pt(12)
                                before_run.font.bold = False
                            
                            # Add the number with bold formatting and larger font
                            num_text = paragraph_text.strip()[start:end]
                            num_run = p.add_run()
                            num_run.text = num_text
                            num_run.font.bold = True
                            num_run.font.size = Pt(14)  # Slightly larger for numbers
                            
                            last_pos = end
                        
                        # Add any remaining text after the last number
                        if last_pos < len(paragraph_text.strip()):
                            after_text = paragraph_text.strip()[last_pos:]
                            after_run = p.add_run()
                            after_run.text = after_text
                            after_run.font.size = Pt(12)
                            after_run.font.bold = False
                    else:
                        # No numbers, just add the text normally
                        p.text = paragraph_text.strip()
                        p.font.bold = False
                        p.font.size = Pt(12)
                    
                    p.space_before = 0
                    p.space_after = 0
                
                # Add other text if present
                if other_text.strip():
                    p = text_frame.add_paragraph()
                    p.text = other_text.strip()
                    p.font.bold = False
                    p.font.size = Pt(12)
                
                # Process images if present - IMPROVED IMAGE POSITIONING
                if has_images:
                    print(f"Processing {len(img_tags)} images in column row")
                    
                    # IMPROVEMENT: Position for the first image - ensure it's below text content
                    # Use a fixed position that's safely below the text
                    img_y = current_y + text_height_actual + Inches(0.3)
                    
                    # Make sure image_y is within the box
                    max_img_y = current_y + box_height - Inches(0.8)  # Leave room at bottom
                    if img_y > max_img_y:
                        img_y = max_img_y - Inches(0.1)  # Emergency adjustment
                    
                    for img_index, img in enumerate(img_tags):
                        try:
                            img_url = img.get('src', '')
                            print(f"Processing image {img_index+1}: {img_url}")
                            
                            if img_url:
                                # Increase timeout to help with connection issues
                                response = requests.get(img_url, stream=True, timeout=15)
                                if response.status_code == 200:
                                    img_bytes = BytesIO(response.content)
                                    
                                    try:
                                        # Get dimensions from image
                                        with PILImage.open(img_bytes) as pil_img:
                                            original_width, original_height = pil_img.size
                                            aspect_ratio = original_width / original_height
                                            
                                            img_bytes.seek(0)  # Reset file pointer
                                            
                                            # Calculate image size - IMPROVED SIZING LOGIC
                                            img_width = None
                                            img_height = None
                                            
                                            # If both width and height specified, use those as starting point
                                            if img.get('width') and img.get('height'):
                                                try:
                                                    width_px = int(img.get('width'))
                                                    height_px = int(img.get('height'))
                                                    
                                                    # Apply minimum sizes
                                                    width_px = max(width_px, 50)  # Minimum 50px
                                                    height_px = max(height_px, 50)  # Minimum 50px
                                                    
                                                    img_width = Inches(width_px / 96)
                                                    img_height = Inches(height_px / 96)
                                                except (ValueError, TypeError):
                                                    # Fall back to calculated dimensions
                                                    img_width = min(Inches(width / 2), Inches(2.5))
                                                    img_height = img_width / aspect_ratio
                                            else:
                                                # No dimensions specified, calculate based on available space
                                                # Use a smaller fraction of column width
                                                img_width = min(width * 0.8, Inches(2.5))
                                                img_height = img_width / aspect_ratio
                                            
                                            # Calculate remaining space within the shape
                                            remaining_height = (current_y + box_height) - img_y - Inches(0.2)
                                            
                                            # Ensure image fits within available height
                                            if img_height > remaining_height and remaining_height > Inches(0.3):
                                                img_height = remaining_height
                                                img_width = img_height * aspect_ratio
                                            
                                            # Ensure image fits within column width
                                            max_width = width - Inches(0.4)
                                            if img_width > max_width:
                                                img_width = max_width
                                                img_height = img_width / aspect_ratio
                                            
                                            # Center the image horizontally
                                            img_x = x_pos + (width - img_width) / 2
                                            
                                            # Final check to ensure reasonable dimensions
                                            if img_width < Inches(0.2) or img_height < Inches(0.2):
                                                # Skip if image would be too small
                                                print(f"Skipping too small image: {img_width} x {img_height}")
                                                continue
                                            
                                            # Only add if we have valid dimensions
                                            if img_width > 0 and img_height > 0:
                                                # IMPROVED: Create a new copy of image data for safety
                                                img_data = BytesIO(img_bytes.getvalue())
                                                
                                                picture = slide.shapes.add_picture(
                                                    img_data, 
                                                    img_x, 
                                                    img_y, 
                                                    width=img_width, 
                                                    height=img_height
                                                )
                                                
                                                img_data.close()  # Close the copy after use
                                                print(f"Added image in column from {img_url} at position: {img_x}, {img_y}, size: {img_width} x {img_height}")
                                            else:
                                                print(f"Invalid image dimensions calculated: {img_width} x {img_height}")
                                    except Exception as img_error:
                                        print(f"Error processing column image: {img_error}")
                                    
                                    img_bytes.close()
                                else:
                                    print(f"Image download failed with status code: {response.status_code}")
                        except Exception as img_error:
                            print(f"Error with column image: {img_error}")
                
                # Update position for next row - ENSURE ADEQUATE SPACING
                current_y += box_height + Inches(0.3)  # Increased spacing between rows
            
            except Exception as row_error:
                print(f"Error processing column row: {row_error}")
                current_y += Inches(0.5)
    
    except Exception as column_error:
        print(f"Error processing column: {column_error}")
    
    return current_y


# Helper function to calculate appropriate box height
def calculate_dynamic_box_height(header_text, paragraph_text, other_text, has_images, image_height):
    """Calculate box height more accurately based on content"""
    # Base height for padding
    padding = Inches(0.4)  # Increased from 0.3
    
    # Calculate text content size
    header_lines = 0
    if header_text.strip():
        # Headers usually need more space
        header_lines = (len(header_text) // 40) + 1
        header_height = Inches(0.25 * header_lines)  # More space for headers
    else:
        header_height = Inches(0)
    
    # Paragraph text
    para_lines = 0
    if paragraph_text.strip():
        para_lines = (len(paragraph_text) // 35) + 1  # Characters per line estimate
        para_height = Inches(0.2 * para_lines)
    else:
        para_height = Inches(0)
    
    # Other text
    other_lines = 0
    if other_text.strip():
        other_lines = (len(other_text) // 40) + 1
        other_height = Inches(0.2 * other_lines)
    else:
        other_height = Inches(0)
    
    # Text height with spacing between elements
    text_spacing = Inches(0.1 * max(1, (header_lines > 0) + (para_lines > 0) + (other_lines > 0) - 1))
    text_height = header_height + para_height + other_height + text_spacing
    
    # Add space for images if present - INCREASED BUFFER
    image_buffer = Inches(0.5) if has_images else Inches(0)
    total_height = text_height + image_buffer + image_height + padding
    
    # IMPROVED: Add minimum height requirement for rows with images
    if has_images:
        min_image_row_height = Inches(1.5)  # Minimum height for rows with images
        total_height = max(total_height, min_image_row_height)
    
    # Ensure minimum reasonable height
    return max(total_height, Inches(0.7))


def estimate_row_height(row):
    """More accurate estimation of row height based on content quantity"""
    # Base height for any row
    height = Inches(0.5)
    
    # Get text content
    text_content = row.get_text().strip()
    text_length = len(text_content)
    
    # Calculate height based on text length with more realistic estimates
    # Assuming approximately 40 characters per line and 0.2 inches per line
    if text_length > 0:
        lines = max(1, text_length // 40)
        text_height = Inches(0.2 * lines)
        height = max(height, text_height)
    
    # Add height for images
    img = row.find('img')
    if img:
        # If height attribute exists, use it
        if img.get('height'):
            try:
                img_height = int(img.get('height')) / 96  # Convert px to inches
                height = max(height, Inches(img_height + 0.4))  # Add margin
            except (ValueError, TypeError):
                height = max(height, Inches(2.0))  # Default if can't parse
        else:
            # Default height for images
            height = max(height, Inches(2.0))
    
    # Add height for tables
    if row.find('table'):
        rows = len(row.find_all('tr'))
        height = max(height, Inches(0.3 * rows + 0.3))  # 0.3 inches per row plus header
    
    # Add height for code blocks
    code_block = row.find('div', class_='code-block') or row.find('pre')
    if code_block:
        code_text = code_block.get_text().strip()
        code_lines = len(code_text.split('\n'))
        height = max(height, Inches(0.2 * code_lines + 0.3))  # 0.2 inches per line
    
    # Handle special elements
    if row.find('ul') or row.find('ol'):
        list_items = len(row.find_all('li'))
        height = max(height, Inches(0.25 * list_items + 0.3))  # 0.25 inches per list item
    
    # Add extra padding to prevent content being cut off
    return height + Inches(0.2)

def add_textbox_relative(slide, top, left, width, height, text, font_size=14, bg_color=None):
    # Optional: add a background shape
    if bg_color:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = RGBColor(200, 200, 200)

    # Add the actual textbox
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = text

    # Format text
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = False

    return top + height + Inches(0.1)  # Return next top position



def process_content(element, text_frame, slide, y_position=None, prs=None, slide_index=0):
    max_y = y_position if y_position is not None else Inches(1.5)
    
    process_headers_with_color(element, text_frame)
    process_paragraphs_with_color(element, text_frame)
    
    text_height = Inches(0.3) * len(text_frame.paragraphs)
    
    if element.find('table'):
        process_table(element.find('table'), text_frame)
    elif element.find('ul') or element.find('ol'):
        process_list(element, text_frame)
    elif element.find(['pre', 'code']) or element.find('div', class_='code-block'):
        process_code_block(element, text_frame)
    
    img = element.find('img')
    if img:
        img_top = max_y + text_height + Inches(0.2)
        
        img_url = img.get('src', '')
        img_alt = img.get('alt', 'Image')
        
        try:
            response = requests.get(img_url, stream=True, timeout=10)
            
            if response.status_code == 200:
                img_bytes = BytesIO(response.content)
                
                try:
                    with PILImage.open(img_bytes) as pil_img:
                        img_width, img_height = pil_img.size
                        aspect_ratio = img_width / img_height
                    
                    img_bytes.seek(0)
                    
                    width_specified = img.get('width')
                    height_specified = img.get('height')
                    
                    img_width = Inches(2.0)
                    img_height = img_width / aspect_ratio
                    
                    if width_specified and height_specified:
                        try:
                            width_px = int(width_specified)
                            height_px = int(height_specified)
                            if width_px > 0 and height_px > 0:
                                img_width = Inches(width_px / 96)
                                img_height = Inches(height_px / 96)
                        except (ValueError, TypeError):
                            pass
                    
                    slide_width = Inches(SLIDE_WIDTH_INCHES)
                    left_position = (slide_width - img_width) / 2
                    
                    if img_width > Inches(6):
                        img_width = Inches(6)
                        img_height = img_width / aspect_ratio
                    
                    picture = slide.shapes.add_picture(
                        img_bytes, 
                        left_position, 
                        img_top, 
                        width=img_width, 
                        height=img_height
                    )
                    
                    max_y = max(max_y, img_top + img_height + Inches(0.2))
                    
                except Exception as img_error:
                    print(f"Error processing image: {img_error}")
                    p = text_frame.add_paragraph()
                    p.text = f"[Image Error: {img_alt}]"
            else:
                p = text_frame.add_paragraph()
                p.text = f"[Image not available: {img_alt}]"
                
        except Exception as request_error:
            print(f"Error downloading image: {request_error}")
            p = text_frame.add_paragraph()
            p.text = f"[Image download error: {img_alt}]"
    
    return max_y




def process_list(element, text_frame):
    """Process HTML lists and add them to the text frame"""
    # First add any text before the list
    text_before = ''
    list_elem = element.find(['ul', 'ol'])
    
    for sibling in list_elem.previous_siblings:
        if isinstance(sibling, str) and sibling.strip():
            text_before += sibling.strip() + ' '
        elif hasattr(sibling, 'get_text'):
            text_before += sibling.get_text().strip() + ' '
            
    if text_before.strip():
        p = text_frame.add_paragraph()
        p.text = text_before.strip()
    
    # Process list items
    is_ordered = list_elem.name == 'ol'
    list_items = list_elem.find_all('li')
    
    for i, item in enumerate(list_items):
        p = text_frame.add_paragraph()
        prefix = f"{i+1}. " if is_ordered else "• "
        p.text = prefix + item.get_text().strip()
        p.level = 1  # Set indentation level
        
        

def process_table(table, text_frame):
    """Process HTML table and add it to the text frame as formatted text"""
    # Add table caption or heading
    p = text_frame.add_paragraph()
    p.text = "[Table]"
    p.font.bold = True
    
    # Process headers
    headers = [th.get_text().strip() for th in table.find_all('th')]
    if headers:
        p = text_frame.add_paragraph()
        p.text = " | ".join(headers)
        p.font.bold = True
        
        # Add separator line
        p = text_frame.add_paragraph()
        p.text = "-" * (sum(len(h) for h in headers) + 3 * (len(headers) - 1))
    
    # Process rows
    for row in table.find_all('tr'):
        cells = [td.get_text().strip() for td in row.find_all('td')]
        if cells:
            p = text_frame.add_paragraph()
            p.text = " | ".join(cells)

def process_code_block(element, text_frame):
    """Process code blocks and add them to the text frame"""
    # Find the code block element
    code_elem = element.find(['pre', 'code']) or element.find('div', class_='code-block')
    
    if not code_elem:
        return
        
    # Add a label
    p = text_frame.add_paragraph()
    p.text = "[Code]"
    p.font.bold = True
    
    # Process code lines
    code_text = code_elem.get_text().strip()
    lines = code_text.split('\n')
    
    for line in lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(9)
def process_image_with_download(element, text_frame, slide, css_rules, y_position=None):
    """Process images with improved error handling to prevent file corruption"""
    img = element.find('img')
    if not img:
        return y_position
    
    # Get image attributes
    img_url = img.get('src', '')
    img_alt = img.get('alt', 'Image')
    
    # Use standard slide dimensions
    slide_width_inches = SLIDE_WIDTH_INCHES
    slide_height_inches = SLIDE_HEIGHT_INCHES
    
    # Calculate content area
    left = Inches(0.5)
    top = y_position if y_position is not None else Inches(1.5)
    
    # Calculate available height on current slide
    available_height = Inches(slide_height_inches - 1.0) - top  # 1.0 inch margin at bottom for safety
    
    # Skip if not enough space
    if available_height < Inches(0.5):
        p = text_frame.add_paragraph()
        p.text = f"[Image: {img_alt} - not enough space]"
        return y_position
    
    try:
        # Download the image with timeout
        response = requests.get(img_url, stream=True, timeout=10)
        
        if response.status_code != 200:
            # Failed to download image
            p = text_frame.add_paragraph()
            p.text = f"[Image: {img_alt} - download failed]"
            p.alignment = PP_ALIGN.CENTER
            return y_position + Inches(0.5)
        
        # Create image from content
        img_bytes = BytesIO(response.content)
        
        try:
            # Try to open the image to validate it
            with PILImage.open(img_bytes) as pil_img:
                img_width, img_height = pil_img.size
                
                # Skip extremely small or zero-dimension images
                if img_width < 10 or img_height < 10:
                    p = text_frame.add_paragraph()
                    p.text = f"[Image: {img_alt} - invalid dimensions]"
                    p.alignment = PP_ALIGN.CENTER
                    return y_position + Inches(0.5)
                
                aspect_ratio = img_width / img_height
            
            # Reset file pointer
            img_bytes.seek(0)
            
            # Get dimensions from HTML
            width_specified = img.get('width')
            height_specified = img.get('height')
            
            # Default dimensions
            width = Inches(6)  # 6 inches wide by default
            height = Inches(6 / aspect_ratio)
            
            # Try to use HTML dimensions if available
            if width_specified:
                try:
                    width_px = int(width_specified)
                    if 10 <= width_px <= 2000:  # Reasonable range check
                        width = Inches(width_px / 96)
                except (ValueError, TypeError):
                    pass  # Keep default width
            
            if height_specified:
                try:
                    height_px = int(height_specified)
                    if 10 <= height_px <= 2000:  # Reasonable range check
                        height = Inches(height_px / 96)
                except (ValueError, TypeError):
                    pass  # Keep calculated height
            
            # Fit to slide width and available height
            max_width = Inches(slide_width_inches - 1.0)  # 0.5 inch margins on each side
            if width > max_width:
                width = max_width
                height = width / aspect_ratio
            
            if height > available_height:
                height = available_height
                width = height * aspect_ratio
            
            # Set minimum dimensions to avoid errors
            width = max(width, Inches(0.1))
            height = max(height, Inches(0.1))
            
            # Create a new BytesIO object to ensure memory is properly managed
            img_data = BytesIO(img_bytes.getvalue())
            img_bytes.close()  # Close the original BytesIO object
            
            # Add image to slide with explicit error handling
            try:
                picture = slide.shapes.add_picture(img_data, left, top, width=width, height=height)
                
                # Close the BytesIO object after adding the picture
                img_data.close()
                
                # Update position for next element
                new_top = top + height + Inches(0.1)
                
                # Add caption if available
                caption = element.find('p', class_='caption')
                if caption and caption.get_text().strip():
                    caption_text = caption.get_text().strip()
                    
                    # Only add caption if there's space
                    if new_top + Inches(0.3) < Inches(slide_height_inches - 0.5):
                        try:
                            caption_box = slide.shapes.add_textbox(
                                left, new_top, width, Inches(0.3)
                            )
                            caption_frame = caption_box.text_frame
                            p = caption_frame.add_paragraph()
                            p.text = caption_text
                            p.font.italic = True
                            p.alignment = PP_ALIGN.CENTER
                            
                            new_top += Inches(0.4)
                        except Exception as caption_error:
                            print(f"Error adding caption: {caption_error}")
                            # Add caption in text frame instead
                            p = text_frame.add_paragraph()
                            p.text = f"Caption: {caption_text}"
                            p.font.italic = True
                
                return new_top
                
            except Exception as picture_error:
                print(f"Error adding picture to slide: {picture_error}")
                # Fallback to text
                p = text_frame.add_paragraph()
                p.text = f"[Image: {img_alt} - failed to add to slide]"
                p.alignment = PP_ALIGN.CENTER
                return y_position + Inches(0.5)
                
        except Exception as pil_error:
            print(f"Error processing image data: {pil_error}")
            # Invalid image data
            p = text_frame.add_paragraph()
            p.text = f"[Image: {img_alt} - invalid image]"
            p.alignment = PP_ALIGN.CENTER
            return y_position + Inches(0.5)
            
    except Exception as request_error:
        print(f"Error downloading image {img_url}: {request_error}")
        # Failed request
        p = text_frame.add_paragraph()
        p.text = f"[Image: {img_alt} - download error]"
        p.alignment = PP_ALIGN.CENTER
        return y_position + Inches(0.5)






def clean_slide_placeholders(slide):
    """Remove or hide any empty placeholders on the slide"""
    for shape in slide.shapes:
        # Check if it's a placeholder
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            try:
                # Try setting it to empty to remove the "Click to add..." text
                if hasattr(shape, 'text'):
                    shape.text = ""
                # Or try to hide it
                if hasattr(shape, 'element') and hasattr(shape.element, 'getparent'):
                    parent = shape.element.getparent()
                    if parent is not None:
                        parent.remove(shape.element)
            except:
                # If we can't modify it, just continue
                pass

def create_html_file_from_string(html_content, filename="temp_html.html"):
    """
    Create a temporary HTML file from a string
    
    Args:
        html_content (str): HTML content as a string
        filename (str): Filename to save the HTML content
        
    Returns:
        str: Path to the created HTML file
    """
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(html_content)
    return filename
def apply_slide_background_color(slide_html, current_slide):
    """Apply background color to the entire slide based on color classes"""
    try:
        # Get the background color from the slide's class
        bg_color = get_color_from_class(slide_html)
        
        # Get the RGB values - RGBColor objects store RGB values directly in rgb attribute
        default_color = RGBColor(255, 255, 255)
        
        # Direct comparison of RGBColor objects
        if bg_color != default_color:
            # Add a background shape that covers the entire slide
            left = top = 0
            width = Inches(SLIDE_WIDTH_INCHES)
            height = Inches(SLIDE_HEIGHT_INCHES)
            
            # Create a rectangle that covers the entire slide
            bg_shape = current_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                left, top, width, height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = bg_color
            
            # No border
            bg_shape.line.width = 0
            
            # Don't try to reorder shapes - this can corrupt the PowerPoint
            print(f"Applied {get_color_name(bg_color)} background to slide")
    except Exception as e:
        # If background color application fails, log it but don't crash
        print(f"Warning: Could not apply slide background color: {e}")

def get_color_name(color):
    """Get a color name from an RGBColor object by comparing values"""
    # Map RGB tuples to color names
    color_map_reverse = {
        (255, 200, 200): "red",
        (200, 200, 255): "blue", 
        (200, 255, 200): "green",
        (255, 255, 200): "yellow",
        (255, 225, 180): "orange",
        (230, 200, 255): "purple",
        (220, 220, 220): "grey",
        (255, 200, 230): "pink",
        (180, 240, 240): "teal"
    }
    
    # Try to find the color by direct comparison
    for rgb_tuple, name in color_map_reverse.items():
        rgb_color = RGBColor(*rgb_tuple)
        if str(color.rgb) == str(rgb_color.rgb):
            return name
    
    # If no match, return generic description
    return "custom"
def add_footer(slide, footer_text="@surveys"):
    """
    Adds a blue footer with white text to the bottom of the slide.
    """
    left = Inches(0)
    top = Inches(SLIDE_HEIGHT_INCHES - FOOTER_HEIGHT_INCHES)
    width = Inches(SLIDE_WIDTH_INCHES)
    height = Inches(FOOTER_HEIGHT_INCHES)

    # Footer background
    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue
    footer_shape.line.fill.background()

    # Footer text box
    text_box = slide.shapes.add_textbox(
        left + Inches(0.5), top + Inches(0.1),
        width - Inches(1.0), height - Inches(0.2)
    )
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White
    p.alignment = PP_ALIGN.RIGHT

def html_from_file_to_pptx(html_file, output_file="presentation.pptx"):
    """
    Process HTML file and convert to PowerPoint
    
    Args:
        html_file (str): Path to HTML file
        output_file (str): Path to save PowerPoint file
    """
    try:
        with open(html_file, 'r', encoding='utf-8') as f:
            html_content = f.read()
            
        # Convert HTML to PowerPoint
        html_to_pptx(html_content, output_file)
        print(f"Successfully converted {html_file} to {output_file}")
        
    except FileNotFoundError:
        print(f"File not found: {html_file}")
        print("Please ensure the HTML file exists or specify the correct path.")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    # Define file paths
    template_file = "placeholder.html"
    json_file = "data.json"
    output_pptx = "presentation.pptx"
    
    # Optional: Define a banner URL (set to None to use default blue banner)
    banner_url = None  # Replace with your actual banner URL
    
    try:
        # Generate PowerPoint using the updated function with banner URL
        generate_ppt_from_json_and_template(template_file, json_file, output_pptx, banner_url)
        
    except FileNotFoundError as e:
        print(f"Error: File not found - {e}")
    except json.JSONDecodeError:
        print(f"Error: '{json_file}' contains invalid JSON.")
    except Exception as e:
        print(f"Error generating PowerPoint: {e}")


    #after adding the the url access to add the banner and increase in font and banner decrease and numbers highlight