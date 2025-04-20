from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
from io import BytesIO
import requests
import re
import html
import sys
import os
import copy

# Standard slide dimensions in inches
SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5

def html_to_pptx(html_content, output_filename="presentation.pptx"):
    """
    Convert HTML to PowerPoint presentation with support for mixed layouts
    
    Args:
        html_content (str): HTML content with slides
        output_filename (str): Output PowerPoint file name
    """
    # Create a new presentation
    prs = Presentation()
    
    # Parse HTML content
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Extract styles from the HTML
    css_rules = extract_css_rules(soup)
    
    # Find all slide divs
    slides = soup.find_all('div', class_='slide')
    
    # Process each slide based on its content
    for slide_index, slide_html in enumerate(slides):
        # Check if this slide has column layout
        left_column = slide_html.find('div', class_='left-column')
        right_column = slide_html.find('div', class_='right-column')
        use_columns_for_slide = bool(left_column or right_column)
        
        if use_columns_for_slide:
            # Process as column layout
            process_column_slide(slide_html, prs, slide_index, css_rules)
        else:
            # Process as standard layout
            process_standard_slide(slide_html, prs, slide_index, css_rules)
    
    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")

def process_standard_slide(slide, prs, slide_index, css_rules):
    """Process a slide with standard layout"""
    # Use a blank slide to avoid placeholders
    slide_layout = prs.slide_layouts[6]  # Blank slide
    current_slide = prs.slides.add_slide(slide_layout)
    
    # Add title manually instead of using placeholder
    title_element = slide.find('h1') or slide.find('h2') or slide.find('p', class_='strong') or slide.find('strong')
    if title_element:
        title_shape = current_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_frame = title_shape.text_frame
        p = title_frame.add_paragraph()
        p.text = title_element.text.strip()
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    # Process the slide content - now passing prs and slide_index
    process_standard_slide_content(slide, current_slide, css_rules, prs, slide_index)
    
    # Clean up any lingering placeholders
    clean_slide_placeholders(current_slide)

def process_column_slide(slide_html, prs, slide_idx, css_rules):
    """Process a slide with column layout and dynamic image handling with overflow support"""
    # Extract the slide title for reuse in continuation slides
    title_element = slide_html.find('h1') or slide_html.find('h2') or slide_html.find('p', class_='strong') or slide_html.find('strong')
    title_text = title_element.get_text().strip() if title_element else f"Slide {slide_idx + 1}"
    
    # Left and Right columns
    left_column = slide_html.find('div', class_='left-column')
    right_column = slide_html.find('div', class_='right-column')
    
    # Column layout setup
    # Calculate dynamic column widths
    margin = Inches(0.5)
    col_spacing = Inches(0.5)
    
    usable_width = Inches(SLIDE_WIDTH_INCHES - 1 - 0.5)  # Total width minus margins
    col_width = (usable_width - col_spacing) / 2  # Equal width for both columns
    
    left_x = margin
    right_x = margin + col_width + col_spacing
    
    # Create the initial slide
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add the title to the first slide
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(SLIDE_WIDTH_INCHES - 1), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    
    start_y = Inches(1.5)  # Start below title
    
    # Process left column with overflow handling
    current_slide = slide
    current_slide_idx = slide_idx
    y_left = start_y
    
    if left_column:
        # Extract all rows from left column for processing
        left_rows = left_column.find_all('div', class_='row')
        row_index = 0
        
        while row_index < len(left_rows):
            # Check remaining space on the slide
            max_y_left = Inches(SLIDE_HEIGHT_INCHES - 0.7)
            
            if y_left >= max_y_left and row_index < len(left_rows):
                # Need to create a continuation slide
                current_slide_idx += 1
                current_slide = prs.slides.add_slide(slide_layout)
                
                # Add a title indicating continuation
                cont_title_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                                               Inches(SLIDE_WIDTH_INCHES - 1), Inches(0.8))
                cont_title_frame = cont_title_box.text_frame
                cont_p = cont_title_frame.add_paragraph()
                cont_p.text = f"{title_text} (Continued)"
                cont_p.font.size = Pt(24)
                cont_p.font.bold = True
                
                # Reset y positions for new slide
                y_left = Inches(1.5)  # Start below title
            
            # Get the current row
            row = left_rows[row_index]
            
            # Estimate row height to check if it will fit
            row_height = estimate_row_height(row)
            
            # Check if this single row fits on current slide
            if y_left + row_height <= max_y_left:
                # Process this row on current slide
                new_y = process_row_in_column(row, current_slide, left_x, y_left, col_width, css_rules)
                y_left = new_y + Inches(0.3) if new_y else y_left + row_height + Inches(0.3)  # Update position with spacing
                row_index += 1  # Move to next row
            else:
                # This row is too big for remaining space, create a new slide
                if row_index < len(left_rows) and y_left > Inches(1.8):  # If not at top of slide already
                    current_slide_idx += 1
                    current_slide = prs.slides.add_slide(slide_layout)
                    
                    # Add a title indicating continuation
                    cont_title_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                                                  Inches(SLIDE_WIDTH_INCHES - 1), Inches(0.8))
                    cont_title_frame = cont_title_box.text_frame
                    cont_p = cont_title_frame.add_paragraph()
                    cont_p.text = f"{title_text} (Continued)"
                    cont_p.font.size = Pt(24)
                    cont_p.font.bold = True
                    
                    # Reset y position for new slide
                    y_left = Inches(1.5)
                else:
                    # Force processing of the row even if it might not fully fit (for very large rows)
                    new_y = process_row_in_column(row, current_slide, left_x, y_left, col_width, css_rules)
                    y_left = new_y + Inches(0.3) if new_y else y_left + row_height + Inches(0.3)
                    row_index += 1
    
    # Reset slide to first slide for right column
    current_slide = slide
    current_slide_idx = slide_idx
    y_right = start_y
    
    if right_column:
        # Extract all rows from right column for processing
        right_rows = right_column.find_all('div', class_='row')
        row_index = 0
        
        while row_index < len(right_rows):
            # Check remaining space on the slide
            max_y_right = Inches(SLIDE_HEIGHT_INCHES - 0.7)
            
            if y_right >= max_y_right and row_index < len(right_rows):
                # Find the corresponding left column slide or create a new one if needed
                slide_offset = current_slide_idx - slide_idx
                if slide_offset + 1 >= len(prs.slides) - slide_idx:
                    # Need a new slide
                    current_slide = prs.slides.add_slide(slide_layout)
                    
                    # Add a title indicating continuation
                    cont_title_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                                                  Inches(SLIDE_WIDTH_INCHES - 1), Inches(0.8))
                    cont_title_frame = cont_title_box.text_frame
                    cont_p = cont_title_frame.add_paragraph()
                    cont_p.text = f"{title_text} (Continued)"
                    cont_p.font.size = Pt(24)
                    cont_p.font.bold = True
                else:
                    # Use existing slide
                    current_slide_idx += 1
                    slide_offset = current_slide_idx - slide_idx
                    current_slide = prs.slides[slide_idx + slide_offset]
                
                # Reset y position for new slide
                y_right = Inches(1.5)  # Start below title
            
            # Get the current row
            row = right_rows[row_index]
            
            # Estimate row height to check if it will fit
            row_height = estimate_row_height(row)
            
            # Check if this single row fits on current slide
            if y_right + row_height <= max_y_right:
                # Process this row on current slide
                new_y = process_row_in_column(row, current_slide, right_x, y_right, col_width, css_rules)
                y_right = new_y + Inches(0.3) if new_y else y_right + row_height + Inches(0.3)  # Update position with spacing
                row_index += 1  # Move to next row
            else:
                # This row is too big for remaining space, create a new slide
                if row_index < len(right_rows) and y_right > Inches(1.8):  # If not at top of slide already
                    current_slide_idx += 1
                    
                    # Check if we need to create a new slide
                    slide_offset = current_slide_idx - slide_idx
                    if slide_offset >= len(prs.slides) - slide_idx:
                        current_slide = prs.slides.add_slide(slide_layout)
                        
                        # Add a title indicating continuation
                        cont_title_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                                                      Inches(SLIDE_WIDTH_INCHES - 1), Inches(0.8))
                        cont_title_frame = cont_title_box.text_frame
                        cont_p = cont_title_frame.add_paragraph()
                        cont_p.text = f"{title_text} (Continued)"
                        cont_p.font.size = Pt(24)
                        cont_p.font.bold = True
                    else:
                        current_slide = prs.slides[slide_idx + slide_offset]
                    
                    # Reset y position for new slide
                    y_right = Inches(1.5)
                else:
                    # Force processing of the row even if it might not fully fit (for very large rows)
                    new_y = process_row_in_column(row, current_slide, right_x, y_right, col_width, css_rules)
                    y_right = new_y + Inches(0.3) if new_y else y_right + row_height + Inches(0.3)
                    row_index += 1

def process_row_in_column(row, slide, x_pos, y_pos, width, css_rules):
    """Process a single row in a column and return the updated y position"""
    # Get images and text content
    img_tags = row.find_all('img')
    has_images = len(img_tags) > 0
    
    # Extract text content properly
    text = ""
    try:
        # Create a temporary copy of the row to work with
        row_copy = BeautifulSoup(str(row), 'html.parser')
        for img_tag in row_copy.find_all('img'):
            img_tag.decompose()  # Remove image tags
        
        # Collect all text from all nodes
        for element in row_copy.descendants:
            if isinstance(element, str) and element.strip():
                text += element.strip() + " "
        
        text = text.strip()
    except Exception as e:
        print(f"Error extracting text: {e}")
        # Fallback to simpler text extraction
        text = row.get_text().strip()
    
    has_text = bool(text)
    
    # Calculate remaining vertical space on slide
    remaining_height = Inches(SLIDE_HEIGHT_INCHES - 0.7) - y_pos
    
    # If both text and images are present, create a unified box
    if has_text and has_images:
        # Calculate space needed for text - ADAPTIVE HEIGHT
        text_length = len(text)
        
        # Calculate appropriate text height based on content length
        if text_length < 100:
            text_height = Inches(0.6)  # Short text
        elif text_length < 250:
            text_height = Inches(1.0)  # Medium text
        elif text_length < 500:
            text_height = Inches(1.5)  # Longer text
        else:
            text_height = Inches(2.0)  # Very long text
        
        # Space for images
        image_space = Inches(1.2)  # Default space for images
        
        # Calculate total box height - ADAPTIVE
        box_height = text_height + image_space + Inches(0.4)  # Text + images + padding
        
        # Ensure it fits in remaining space, or adjust to fit
        if box_height > remaining_height:
            # If there's reasonable space, use what we have
            if remaining_height > Inches(1.5):
                box_height = remaining_height - Inches(0.1)
                text_height = box_height - image_space - Inches(0.4)  # Adjust text height to fit
        
        # Create the unified box (background shape)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            x_pos, y_pos, 
            width, box_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = get_color_from_class(row)
        bg_shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Add text at the top of the box
        text_box = slide.shapes.add_textbox(
            x_pos + Inches(0.2), 
            y_pos + Inches(0.2), 
            width - Inches(0.4), 
            text_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Split text into paragraphs if very long
        if text_length > 300:
            # Try to split at sentences
            sentences = re.split(r'(?<=[.!?])\s+', text)
            
            # Add first sentence
            p = text_frame.add_paragraph()
            p.text = sentences[0]
            p.font.size = Pt(11)  # Slightly smaller font for long text
            
            # Add remaining sentences as separate paragraphs
            for sentence in sentences[1:]:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence
                    p.font.size = Pt(11)
        else:
            # Add as a single paragraph
            p = text_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(12)
        
        # Starting Y position for the first image
        img_y = y_pos + text_height + Inches(0.2)
        
        # Process first image only (safer)
        if img_tags and img_y + Inches(1.0) < y_pos + box_height:
            try:
                img = img_tags[0]  # Just process the first image
                img_url = img.get('src', '')
                if img_url:
                    response = requests.get(img_url, stream=True, timeout=5)
                    if response.status_code == 200:
                        img_bytes = BytesIO(response.content)
                        
                        # Get image dimensions with aspect ratio
                        try:
                            with PILImage.open(img_bytes) as pil_img:
                                aspect_ratio = pil_img.width / pil_img.height
                            
                            img_bytes.seek(0)  # Reset file pointer
                            
                            # Calculate image dimensions based on available space
                            available_width = width - Inches(0.4)
                            available_height = box_height - text_height - Inches(0.4)
                            
                            # Default dimensions
                            img_width = min(Inches(2.0), available_width)
                            img_height = img_width / aspect_ratio
                            
                            # Adjust if height is too large
                            if img_height > available_height:
                                img_height = available_height
                                img_width = img_height * aspect_ratio
                                
                                # Ensure width isn't too large
                                if img_width > available_width:
                                    img_width = available_width
                                    img_height = img_width / aspect_ratio
                            
                            # Calculate centered position for image
                            img_x = x_pos + (width - img_width) / 2
                            
                            # Create picture with proper sizing
                            picture = slide.shapes.add_picture(
                                img_bytes, 
                                img_x, 
                                img_y, 
                                width=img_width, 
                                height=img_height
                            )
                        except Exception as img_error:
                            print(f"Error calculating image dimensions: {img_error}")
                            # Fallback to fixed size if aspect ratio calculation fails
                            img_width = min(Inches(2.0), width - Inches(0.4))
                            img_height = Inches(1.0)
                            
                            # Calculate centered position for image
                            img_x = x_pos + (width - img_width) / 2
                            
                            picture = slide.shapes.add_picture(
                                img_bytes, 
                                img_x, 
                                img_y, 
                                width=img_width, 
                                height=img_height
                            )
                        
                        img_bytes.close()
            except Exception as img_error:
                print(f"Error with image: {img_error}")
        
        # Return updated position
        return y_pos + box_height
    
    # Handle text-only content with ADAPTIVE HEIGHT
    elif has_text:
        # Calculate appropriate text height based on content length
        text_length = len(text)
        
        if text_length < 100:
            text_height = Inches(0.6)  # Short text
        elif text_length < 250:
            text_height = Inches(1.0)  # Medium text
        elif text_length < 500:
            text_height = Inches(1.5)  # Longer text
        else:
            text_height = Inches(2.0)  # Very long text
        
        # Ensure it fits in remaining space
        if text_height > remaining_height - Inches(0.2):
            text_height = remaining_height - Inches(0.2)
        
        # Create textbox with background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            x_pos, y_pos, 
            width, text_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = get_color_from_class(row) 
        shape.line.color.rgb = RGBColor(200, 200, 200)

        # Add the text
        textbox = slide.shapes.add_textbox(
            x_pos + Inches(0.1), 
            y_pos + Inches(0.1), 
            width - Inches(0.2), 
            text_height - Inches(0.2)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Split text into paragraphs if very long
        if text_length > 300:
            # Try to split at sentences
            sentences = re.split(r'(?<=[.!?])\s+', text)
            
            # Add first sentence to first paragraph
            text_frame.text = sentences[0]
            
            # Add remaining sentences as separate paragraphs
            for sentence in sentences[1:]:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence
                    p.font.size = Pt(11)
        else:
            # Add as a single paragraph
            text_frame.text = text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
        
        # Return updated position
        return y_pos + text_height
    
    # Handle image-only content with adaptive sizing
    elif has_images:
        # Process first image only (safer)
        if img_tags:
            try:
                img = img_tags[0]  # Just process the first image
                img_url = img.get('src', '')
                if img_url:
                    response = requests.get(img_url, stream=True, timeout=5)
                    if response.status_code == 200:
                        img_bytes = BytesIO(response.content)
                        
                        # Get image dimensions with aspect ratio
                        try:
                            with PILImage.open(img_bytes) as pil_img:
                                aspect_ratio = pil_img.width / pil_img.height
                            
                            img_bytes.seek(0)  # Reset file pointer
                            
                            # Adaptive sizing based on available space
                            img_width = min(Inches(3.0), width - Inches(0.4))
                            img_height = img_width / aspect_ratio
                            
                            # Make sure it fits
                            if img_height > remaining_height - Inches(0.3):
                                img_height = remaining_height - Inches(0.3)
                                img_width = img_height * aspect_ratio
                                
                                # Ensure width isn't too large
                                if img_width > width - Inches(0.4):
                                    img_width = width - Inches(0.4)
                                    img_height = img_width / aspect_ratio
                            
                            # Calculate centered position for image
                            img_x = x_pos + (width - img_width) / 2
                            
                            picture = slide.shapes.add_picture(
                                img_bytes, 
                                img_x, 
                                y_pos, 
                                width=img_width, 
                                height=img_height
                            )
                            
                            img_bytes.close()
                            return y_pos + img_height
                        except Exception as img_error:
                            print(f"Error calculating image dimensions: {img_error}")
                            # Fallback to fixed dimensions
                            img_width = min(Inches(2.5), width - Inches(0.4))
                            img_height = Inches(2.0)
                            
                            # Calculate centered position for image
                            img_x = x_pos + (width - img_width) / 2
                            
                            picture = slide.shapes.add_picture(
                                img_bytes, 
                                img_x, 
                                y_pos, 
                                width=img_width, 
                                height=img_height
                            )
                            
                            img_bytes.close()
                            return y_pos + img_height
            except Exception as img_error:
                print(f"Error with image: {img_error}")
    
    # Return the original position if no content was processed
    return y_pos

def process_column_content(column, slide, x_pos, y_pos, width, css_rules=None, slide_index=0, prs=None):
    """Process content within a column with properly aligned non-overlapping unified boxes with adaptive heights"""
    current_y = y_pos
    
    try:
        # Process each row in the column
        for row in column.find_all('div', class_='row'):
            try:
                # Calculate remaining vertical space on slide
                remaining_height = Inches(SLIDE_HEIGHT_INCHES - 1.0) - current_y
                
                # Skip if not enough space left on slide
                if remaining_height < Inches(0.5):
                    # If prs is provided, create a continuation slide
                    if prs:
                        # Create a new slide for overflow content
                        slide_layout = prs.slide_layouts[6]  # Blank slide
                        new_slide = prs.slides.add_slide(slide_layout)
                        
                        # Add a title indicating continuation
                        title_element = column.parent.find('h1') or column.parent.find('h2') 
                        title_text = title_element.get_text().strip() if title_element else f"Slide {slide_index + 1}"
                        
                        title_box = new_slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.3), Inches(SLIDE_WIDTH_INCHES - 1), Inches(0.8)
                        )
                        title_frame = title_box.text_frame
                        p = title_frame.add_paragraph()
                        p.text = f"{title_text} (Continued)"
                        p.font.size = Pt(24)
                        p.font.bold = True
                        
                        # Create a new column instance for the remaining rows
                        remaining_column = copy.copy(column)
                        # Keep only the rows we haven't processed yet
                        processed_rows = column.find_all('div', class_='row')[:column.find_all('div', class_='row').index(row)]
                        for r in processed_rows:
                            r.decompose()
                        
                        # Process the remaining rows on the new slide
                        process_column_content(remaining_column, new_slide, x_pos, Inches(1.5), width, css_rules, slide_index + 1, prs)
                        return current_y
                    else:
                        break
                
                # Get images and text content
                img_tags = row.find_all('img')
                has_images = len(img_tags) > 0
                
                # Extract text content properly to include both tag content and direct text
                text = ""
                try:
                    # Create a temporary copy of the row to work with
                    row_copy = BeautifulSoup(str(row), 'html.parser')
                    for img_tag in row_copy.find_all('img'):
                        img_tag.decompose()  # Remove image tags
                    
                    # Collect all text from all nodes
                    for element in row_copy.descendants:
                        if isinstance(element, str) and element.strip():
                            text += element.strip() + " "
                    
                    text = text.strip()
                except Exception as e:
                    print(f"Error extracting text: {e}")
                    # Fallback to simpler text extraction
                    text = row.get_text().strip()
                
                has_text = bool(text)
                
                # If both text and images are present, create a unified box
                if has_text and has_images:
                    try:
                        # Calculate space needed for text - ADAPTIVE HEIGHT
                        # Estimate how many lines the text will take
                        text_length = len(text)
                        text_words = len(text.split())
                        
                        # Calculate appropriate text height based on content length
                        if text_length < 100:
                            text_height = Inches(0.6)  # Short text
                        elif text_length < 250:
                            text_height = Inches(1.0)  # Medium text
                        elif text_length < 500:
                            text_height = Inches(1.5)  # Longer text
                        else:
                            text_height = Inches(2.0)  # Very long text
                        
                        # Space for images
                        image_space = Inches(1.2)  # Default space for images
                        
                        # Calculate total box height - ADAPTIVE
                        box_height = text_height + image_space + Inches(0.4)  # Text + images + padding
                        
                        # Ensure it fits in remaining space
                        if box_height > remaining_height:
                            # If there's reasonable space, use what we have
                            if remaining_height > Inches(1.5):
                                box_height = remaining_height - Inches(0.1)
                                text_height = box_height - image_space - Inches(0.4)  # Adjust text height to fit
                            else:
                                # Skip this content if not enough space - advance a little and continue
                                current_y += Inches(0.2)
                                continue
                        
                        # Create the unified box (background shape)
                        bg_shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE, 
                            x_pos, current_y, 
                            width, box_height
                        )
                        bg_shape.fill.solid()
                        bg_shape.fill.fore_color.rgb = get_color_from_class(row)
                        bg_shape.line.color.rgb = RGBColor(200, 200, 200)
                        
                        # Add text at the top of the box
                        text_box = slide.shapes.add_textbox(
                            x_pos + Inches(0.2), 
                            current_y + Inches(0.2), 
                            width - Inches(0.4), 
                            text_height
                        )
                        text_frame = text_box.text_frame
                        text_frame.word_wrap = True
                        
                        # Split text into paragraphs if very long
                        if text_length > 300:
                            # Try to split at sentences
                            sentences = re.split(r'(?<=[.!?])\s+', text)
                            
                            # Add first sentence
                            p = text_frame.add_paragraph()
                            p.text = sentences[0]
                            p.font.size = Pt(11)  # Slightly smaller font for long text
                            
                            # Add remaining sentences as separate paragraphs
                            for sentence in sentences[1:]:
                                if sentence.strip():
                                    p = text_frame.add_paragraph()
                                    p.text = sentence
                                    p.font.size = Pt(11)
                        else:
                            # Add as a single paragraph
                            p = text_frame.add_paragraph()
                            p.text = text
                            p.font.size = Pt(12)
                        
                        # Starting Y position for the first image
                        img_y = current_y + text_height + Inches(0.2)
                        
                        # Process first image only (safer)
                        if img_tags and img_y + Inches(1.0) < current_y + box_height:
                            try:
                                img = img_tags[0]  # Just process the first image
                                img_url = img.get('src', '')
                                if img_url:
                                    response = requests.get(img_url, stream=True, timeout=5)
                                    if response.status_code == 200:
                                        img_bytes = BytesIO(response.content)
                                        
                                        # Get image dimensions with aspect ratio
                                        try:
                                            with PILImage.open(img_bytes) as pil_img:
                                                aspect_ratio = pil_img.width / pil_img.height
                                            
                                            img_bytes.seek(0)  # Reset file pointer
                                            
                                            # Calculate image dimensions based on available space
                                            available_width = width - Inches(0.4)
                                            available_height = box_height - text_height - Inches(0.4)
                                            
                                            # Default dimensions
                                            img_width = min(Inches(2.0), available_width)
                                            img_height = img_width / aspect_ratio
                                            
                                            # Adjust if height is too large
                                            if img_height > available_height:
                                                img_height = available_height
                                                img_width = img_height * aspect_ratio
                                                
                                                # Ensure width isn't too large
                                                if img_width > available_width:
                                                    img_width = available_width
                                                    img_height = img_width / aspect_ratio
                                            
                                            # Center the image horizontally
                                            img_x = x_pos + (width - img_width) / 2
                                            
                                            # Create picture with proper sizing
                                            picture = slide.shapes.add_picture(
                                                img_bytes, 
                                                img_x, 
                                                img_y, 
                                                width=img_width, 
                                                height=img_height
                                            )
                                        except:
                                            # Fallback to fixed size if aspect ratio calculation fails
                                            img_width = min(Inches(2.0), width - Inches(0.4))
                                            img_height = Inches(1.0)
                                            
                                            # Center the image horizontally
                                            img_x = x_pos + (width - img_width) / 2
                                            
                                            picture = slide.shapes.add_picture(
                                                img_bytes, 
                                                img_x, 
                                                img_y, 
                                                width=img_width, 
                                                height=img_height
                                            )
                                        
                                        img_bytes.close()
                            except Exception as img_error:
                                print(f"Error with image: {img_error}")
                        
                        # Update position for next row
                        current_y += box_height + Inches(0.3)
                        
                    except Exception as unified_error:
                        print(f"Error creating unified box: {unified_error}")
                        # Skip to next row on error
                        current_y += Inches(0.5)
    
    except Exception as column_error:
        print(f"Error processing column: {column_error}")
    
    return current_y


def process_standard_slide_content(slide_html, current_slide, css_rules, prs=None, slide_index=0):
    """Process content for a standard slide layout with better overflow handling"""
    # Track vertical position for adding content
    current_y = Inches(1.5)  # Start after title
    
    # Get overall text length to determine if we might need special handling
    full_text = slide_html.get_text().strip()
    
    # If the entire content is very long, handle it as overflow text
    if len(full_text) > 1200 and prs:  # Lower threshold for better fit
        content_shape = current_slide.shapes.add_textbox(
            Inches(0.5), current_y, Inches(9), Inches(5)
        )
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        
        # Handle as overflow text
        handle_text_overflow(full_text, content_frame, current_slide, slide_index, prs)
        return
    
    # Find and process all row divs
    rows = slide_html.find_all('div', class_='row')
    
    # If no rows are found, process the slide content directly
    if not rows:
        content_shape = current_slide.shapes.add_textbox(
            Inches(0.5), current_y, Inches(9), Inches(5)
        )
        content_frame = content_shape.text_frame
        process_content(slide_html, content_frame, current_slide, css_rules, current_y, prs, slide_index)
    else:
        # Process each row with better spacing management
        for row_index, row in enumerate(rows):
            # Check remaining space
            remaining_height = Inches(SLIDE_HEIGHT_INCHES - 0.7) - current_y
            if remaining_height < Inches(1.0):
                # Not enough space for meaningful content
                # Create a new slide for remaining content
                if prs and row_index < len(rows):
                    next_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    # Add a title indicating continuation
                    title_shape = next_slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
                    )
                    title_frame = title_shape.text_frame
                    p = title_frame.add_paragraph()
                    p.text = f"Continued from Slide {slide_index+1}"
                    p.font.italic = True
                    p.font.bold = True
                    p.font.size = Pt(18)
                    
                    # Process remaining rows on new slide
                    next_y = Inches(1.5)
                    for next_row in rows[row_index:]:
                        row_height = estimate_row_height(next_row)
                        
                        text_shape = next_slide.shapes.add_textbox(
                            Inches(0.5), next_y, Inches(9), row_height
                        )
                        text_frame = text_shape.text_frame
                        
                        new_y = process_content(next_row, text_frame, next_slide, css_rules, next_y, prs, slide_index+1)
                        
                        next_y = max(next_y + row_height, new_y) + Inches(0.3) if new_y else next_y + row_height + Inches(0.3)
                        
                        # Check if we're running out of space on this slide too
                        if next_y > Inches(SLIDE_HEIGHT_INCHES - 0.7):
                            break
                    
                    # No need to process more rows on the original slide
                    break
            
            # Estimate row height with more conservative calculation
            row_height = estimate_row_height(row)
            
            # Create a text frame for this row
            text_shape = current_slide.shapes.add_textbox(
                Inches(0.5), current_y, Inches(9), row_height
            )
            text_frame = text_shape.text_frame
            
            # Process the content of the row
            new_y = process_content(row, text_frame, current_slide, css_rules, current_y, prs, slide_index)
            
            # Update the vertical position for the next row with MORE SPACE
            current_y = max(current_y + row_height, new_y) + Inches(0.3) if new_y else current_y + row_height + Inches(0.3)


def handle_text_overflow(text, text_frame, slide, current_slide_index, prs):
    """Break long text content across multiple slides with improved handling"""
    # Calculate approximately how much text fits on one slide - MORE CONSERVATIVE
    chars_per_slide = 800  # Reduced from 1500 for better fit
    
    if len(text) > chars_per_slide:
        # Find a good break point - end of sentence or paragraph
        break_point = chars_per_slide
        while break_point > chars_per_slide / 2:
            if text[break_point] in '.!?' and (break_point + 1 >= len(text) or text[break_point + 1] in ' \n\r\t'):
                break_point += 1  # Include the punctuation
                break
            elif text[break_point] in ' \n\r\t' and (break_point > 0 and text[break_point - 1] in '.!?'):
                break
            break_point -= 1
        
        if break_point <= chars_per_slide / 2:
            # If no good break found, find a word boundary
            break_point = chars_per_slide
            while break_point < len(text) and text[break_point] not in ' \n\r\t':
                break_point -= 1
            if break_point <= chars_per_slide / 2:
                break_point = chars_per_slide  # Fall back to hard break
        
        # Add text that fits to current slide
        p = text_frame.add_paragraph()
        p.text = text[:break_point].strip()
        
        # Create a new slide for remaining text with BETTER FORMATTING
        next_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        
        # Add a title indicating continuation
        title_shape = next_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        p = title_frame.add_paragraph()
        p.text = f"Continued from Slide {current_slide_index+1}"
        p.font.italic = True
        p.font.bold = True
        p.font.size = Pt(18)
        
        # Add the content with better positioning
        next_text_shape = next_slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
        )
        next_text_frame = next_text_shape.text_frame
        next_text_frame.word_wrap = True
        
        # Recursively handle remaining text
        remaining_text = text[break_point:].strip()
        handle_text_overflow(remaining_text, next_text_frame, next_slide, 
                            current_slide_index+1, prs)
        
        return True
    else:
        # Just add the text as a paragraph - no overflow
        p = text_frame.add_paragraph()
        p.text = text
        return False


def process_content(element, text_frame, slide, css_rules, y_position=None, prs=None, slide_index=0):
    """Process HTML content and add it to a PowerPoint slide"""
    # Keep track of the vertical position
    max_y = y_position if y_position is not None else Inches(1.5)
    
    # Handle different content types appropriately
    if element.find('table'):
        process_table(element.find('table'), text_frame, css_rules)
    elif element.find('ul') or element.find('ol'):
        process_list(element, text_frame, css_rules)
    elif element.find(['pre', 'code']) or element.find('div', class_='code-block'):
        process_code_block(element, text_frame, css_rules)
    elif element.find('img'):
        new_y = process_image_with_download(element, text_frame, slide, css_rules, y_position)
        max_y = max(max_y, new_y) if new_y else max_y
    else:
        # Process text content
        process_text_content(element, text_frame, css_rules, slide, prs, slide_index)
    
    return max_y


def process_text_content(element, text_frame, css_rules, slide=None, prs=None, slide_index=0):
    """Process text content and add it to the text frame with simple reliable handling"""
    # Enable word wrap for the text frame
    text_frame.word_wrap = True
    
    # Extract all text with a simpler approach
    all_text = element.get_text().strip()
    
    # Simplify by just adding all text to a single paragraph
    if all_text:
        # For very long text, use overflow handling
        if slide and prs and len(all_text) > 1000:
            # Use simple overflow handler
            chars_per_slide = 1000
            first_part = all_text[:chars_per_slide] + "..."
            
            p = text_frame.add_paragraph()
            p.text = first_part
            
            # Create a new slide for remaining text
            next_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Add title to continuation slide
            title_shape = next_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
            )
            title_frame = title_shape.text_frame
            p = title_frame.add_paragraph()
            p.text = f"Continued from previous slide"
            p.font.italic = True
            p.font.bold = True
            
            # Add content to continuation slide
            next_text_shape = next_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
            )
            next_text_frame = next_text_shape.text_frame
            next_text_frame.word_wrap = True
            
            p = next_text_frame.add_paragraph()
            p.text = all_text[chars_per_slide:]
        else:
            # Just add text directly
            p = text_frame.add_paragraph()
            p.text = all_text
            
            # Apply basic formatting if needed
            if element.name in ['h3', 'h4']:
                p.font.bold = True
                size_map = {'h3': 18, 'h4': 16}
                p.font.size = Pt(size_map.get(element.name, 14))


def process_list(element, text_frame, css_rules):
    """Process HTML lists and add them to the text frame"""
    # First add any text before the list
    text_before = ''
    list_elem = element.find(['ul', 'ol'])
    
    for sibling in list_elem.previous_siblings:
        if isinstance(sibling, str) and sibling.strip():
            text_before += sibling.strip() + ' '
        elif hasattr(sibling, 'get_text'):
            text_before += sibling.get_text().strip() + ' '
            
    if text_before.strip():
        p = text_frame.add_paragraph()
        p.text = text_before.strip()
    
    # Process list items
    is_ordered = list_elem.name == 'ol'
    list_items = list_elem.find_all('li')
    
    for i, item in enumerate(list_items):
        p = text_frame.add_paragraph()
        prefix = f"{i+1}. " if is_ordered else "• "
        p.text = prefix + item.get_text().strip()
        p.level = 1  # Set indentation level
        
        apply_css_to_paragraph(p, item, css_rules)


def process_table(table, text_frame, css_rules):
    """Process HTML table and add it to the text frame as formatted text"""
    # Add table caption or heading
    p = text_frame.add_paragraph()
    p.text = "[Table]"
    p.font.bold = True
    
    # Process headers
    headers = [th.get_text().strip() for th in table.find_all('th')]
    if headers:
        p = text_frame.add_paragraph()
        p.text = " | ".join(headers)
        p.font.bold = True
        
        # Add separator line
        p = text_frame.add_paragraph()
        p.text = "-" * (sum(len(h) for h in headers) + 3 * (len(headers) - 1))
    
    # Process rows
    for row in table.find_all('tr'):
        cells = [td.get_text().strip() for td in row.find_all('td')]
        if cells:
            p = text_frame.add_paragraph()
            p.text = " | ".join(cells)


def process_code_block(element, text_frame, css_rules):
    """Process code blocks and add them to the text frame"""
    # Find the code block element
    code_elem = element.find(['pre', 'code']) or element.find('div', class_='code-block')
    
    if not code_elem:
        return
        
    # Add a label
    p = text_frame.add_paragraph()
    p.text = "[Code]"
    p.font.bold = True
    
    # Process code lines
    code_text = code_elem.get_text().strip()
    lines = code_text.split('\n')
    
    for line in lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(9)


def process_image_with_download(element, text_frame, slide, css_rules, y_position=None):
    """Process images with improved error handling to prevent file corruption"""
    img = element.find('img')
    if not img:
        return y_position
    
    # Get image attributes
    img_url = img.get('src', '')
    img_alt = img.get('alt', 'Image')
    
    # Use standard slide dimensions
    slide_width_inches = SLIDE_WIDTH_INCHES
    slide_height_inches = SLIDE_HEIGHT_INCHES
    
    # Calculate content area
    left = Inches(0.5)
    top = y_position if y_position is not None else Inches(1.5)
    
    # Calculate available height on current slide
    available_height = Inches(slide_height_inches - 1.0) - top  # 1.0 inch margin at bottom for safety
    
    # Skip if not enough space
    if available_height < Inches(0.5):
        p = text_frame.add_paragraph()
        p.text = f"[Image: {img_alt} - not enough space]"
        return y_position
    
    try:
        # Download the image with timeout
        response = requests.get(img_url, stream=True, timeout=10)
        
        if response.status_code != 200:
            # Failed to download image
            p = text_frame.add_paragraph()
            p.text = f"[Image: {img_alt} - download failed]"
            p.alignment = PP_ALIGN.CENTER
            return y_position + Inches(0.5)
        
        # Create image from content
        img_bytes = BytesIO(response.content)
        
        try:
            # Try to open the image to validate it
            with PILImage.open(img_bytes) as pil_img:
                img_width, img_height = pil_img.size
                
                # Skip extremely small or zero-dimension images
                if img_width < 10 or img_height < 10:
                    p = text_frame.add_paragraph()
                    p.text = f"[Image: {img_alt} - invalid dimensions]"
                    p.alignment = PP_ALIGN.CENTER
                    return y_position + Inches(0.5)
                
                aspect_ratio = img_width / img_height
            
            # Reset file pointer
            img_bytes.seek(0)
            
            # Get dimensions from HTML
            width_specified = img.get('width')
            height_specified = img.get('height')
            
            # Default dimensions
            width = Inches(6)  # 6 inches wide by default
            height = Inches(6 / aspect_ratio)
            
            # Try to use HTML dimensions if available
            if width_specified:
                try:
                    width_px = int(width_specified)
                    if 10 <= width_px <= 2000:  # Reasonable range check
                        width = Inches(width_px / 96)
                except (ValueError, TypeError):
                    pass  # Keep default width
                
                # Handle text-only content with ADAPTIVE HEIGHT
                elif has_text:
                    try:
                        # Calculate appropriate text height based on content length
                        text_length = len(text)
                        
                        if text_length < 100:
                            text_height = Inches(0.6)  # Short text
                        elif text_length < 250:
                            text_height = Inches(1.0)  # Medium text
                        elif text_length < 500:
                            text_height = Inches(1.5)  # Longer text
                        else:
                            text_height = Inches(2.0)  # Very long text
                        
                        # Ensure it fits in remaining space
                        if text_height > remaining_height - Inches(0.2):
                            text_height = remaining_height - Inches(0.2)
                        
                        # Create textbox with background
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE, 
                            x_pos, current_y, 
                            width, text_height
                        )
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = get_color_from_class(row) 
                        shape.line.color.rgb = RGBColor(200, 200, 200)

                        # Add the text
                        textbox = slide.shapes.add_textbox(
                            x_pos + Inches(0.1), 
                            current_y + Inches(0.1), 
                            width - Inches(0.2), 
                            text_height - Inches(0.2)
                        )
                        text_frame = textbox.text_frame
                        text_frame.word_wrap = True
                        
                        # Split text into paragraphs if very long
                        if text_length > 300:
                            # Try to split at sentences
                            sentences = re.split(r'(?<=[.!?])\s+', text)
                            
                            # Add first sentence to first paragraph
                            text_frame.text = sentences[0]
                            
                            # Add remaining sentences as separate paragraphs
                            for sentence in sentences[1:]:
                                if sentence.strip():
                                    p = text_frame.add_paragraph()
                                    p.text = sentence
                                    p.font.size = Pt(11)
                        else:
                            # Add as a single paragraph
                            text_frame.text = text
                            paragraph = text_frame.paragraphs[0]
                            paragraph.font.size = Pt(12)
                        
                        # Update position
                        current_y += text_height + Inches(0.3)
                        
                    except Exception as text_error:
                        print(f"Error processing text-only content: {text_error}")
                        current_y += Inches(0.5)
                
                # Handle image-only content with adaptive sizing
                elif has_images:
                    try:
                        # Process first image only (safer)
                        if img_tags:
                            try:
                                img = img_tags[0]  # Just process the first image
                                img_url = img.get('src', '')
                                if img_url:
                                    response = requests.get(img_url, stream=True, timeout=5)
                                    if response.status_code == 200:
                                        img_bytes = BytesIO(response.content)
                                        
                                        # Get image dimensions with aspect ratio
                                        try:
                                            with PILImage.open(img_bytes) as pil_img:
                                                aspect_ratio = pil_img.width / pil_img.height
                                            
                                            img_bytes.seek(0)  # Reset file pointer
                                            
                                            # Adaptive sizing based on available space
                                            img_width = min(Inches(3.0), width - Inches(0.4))
                                            img_height = img_width / aspect_ratio
                                            
                                            # Make sure it fits
                                            if img_height > remaining_height - Inches(0.3):
                                                img_height = remaining_height - Inches(0.3)
                                                img_width = img_height * aspect_ratio
                                                
                                                # Ensure width isn't too large
                                                if img_width > width - Inches(0.4):
                                                    img_width = width - Inches(0.4)
                                                    img_height = img_width / aspect_ratio
                                            
                                            # Center the image horizontally
                                            img_x = x_pos + (width - img_width) / 2
                                            
                                            picture = slide.shapes.add_picture(
                                                img_bytes, 
                                                img_x, 
                                                current_y, 
                                                width=img_width, 
                                                height=img_height
                                            )
                                            
                                            img_bytes.close()
                                            current_y += img_height + Inches(0.3)
                                        except Exception as img_calc_error:
                                            print(f"Error calculating image dimensions: {img_calc_error}")
                                            # Fallback to fixed dimensions
                                            img_width = min(Inches(2.5), width - Inches(0.4))
                                            img_height = Inches(2.0)
                                            
                                            # Center the image horizontally
                                            img_x = x_pos + (width - img_width) / 2
                                            
                                            picture = slide.shapes.add_picture(
                                                img_bytes, 
                                                img_x, 
                                                current_y, 
                                                width=img_width, 
                                                height=img_height
                                            )
                                            
                                            img_bytes.close()
                                            current_y += img_height + Inches(0.3)
                            except Exception as img_error:
                                print(f"Error with image: {img_error}")
                                current_y += Inches(0.5)
                    except Exception as img_section_error:
                        print(f"Error in image section: {img_section_error}")
                        current_y += Inches(0.5)
                
                # Add spacing between rows
                current_y += Inches(0.15)
                
            except Exception as row_error:
                print(f"Error processing row: {row_error}")
                current_y += Inches(0.5)